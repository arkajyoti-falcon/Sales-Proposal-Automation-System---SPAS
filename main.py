from __future__ import annotations
import base64
import io
import os
import re
import json
from PIL import Image
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass
from datetime import datetime
from typing import Dict, List, Optional, Tuple
from typing import Optional
import fitz  # PyMuPDF
import streamlit as st
import time
import glob
import pathlib
import tempfile
from docx import Document
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from docxcompose.composer import Composer
from groq import Groq
from dotenv import load_dotenv
import streamlit.components.v1 as components
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage, ImageChops, ImageOps
from docx.enum.style import WD_STYLE_TYPE

# Optional HTTP client for Mermaid→PNG rendering via Kroki
try:
    import requests  # noqa
    HAS_REQUESTS = True
except Exception:
    HAS_REQUESTS = False

# Optional OCR support (not strictly required)
try:
    from PIL import Image as PILImage
    import pytesseract
    HAS_TESS = True
except Exception:
    HAS_TESS = False

# Optional SVG→PNG rasterizer (only used if present)
try:
    import cairosvg
    HAS_CAIROSVG = True
except Exception:
    HAS_CAIROSVG = False

import zlib  # used for mermaid.ink encoding

# Optional PPTX export (editable flowchart)
try:
    from pptx import Presentation
    from pptx.util import Inches as PPTInches, Pt as PPTPt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.dml.color import RGBColor as PPTXRGB
    HAS_PPTX = True
except Exception:
    HAS_PPTX = False

# Optional Selenium (fallback Mermaid→PNG and optional draw.io helper)
try:
    from selenium import webdriver
    from selenium.webdriver.chrome.service import Service as ChromeService
    from selenium.webdriver.common.by import By
    from selenium.webdriver.support.ui import WebDriverWait
    from selenium.webdriver.support import expected_conditions as EC
    from webdriver_manager.chrome import ChromeDriverManager
    HAS_SELENIUM = True
except Exception:
    HAS_SELENIUM = False

load_dotenv()

# ============================================================================
# Configuration & Clients
# ============================================================================

# Read API keys from environment; fall back to provided values to preserve behavior
# Groq client (single instance used everywhere)
try:
    groq_client  # noqa: F401
except NameError:
    _key = os.getenv("GROQ_API_KEY", None)
    groq_client = Groq(api_key=_key)

# ============================================================================
# Prompt (you can paste your big prompts here)
# ============================================================================
PROMPT_COVER_LETTER = """
You are an AI assistant working as a professional proposal writer at Falcon Autotech. You are an expert in drafting formal, client-specific techno-commercial cover letters for proposals. Your role is to generate well-structured, personalized cover letters that follow Falcon’s business communication style, maintain a professional and respectful tone, and clearly demonstrate Falcon’s commitment, expertise, and partnership approach to clients.
Generate a formal techno-commercial COVER LETTER for a proposal. 
The writing style MUST be indistinguishable from natural human writing. The text should read as if drafted by an experienced professional, not an AI system. Use clear, simple, and natural language with varied sentence lengths and structures. Avoid generic phrases, repetitive patterns, or mechanical tone. Ensure that the output flows smoothly, conveys intent naturally, and would not be detected as machine-generated. The content should feel thoughtful, context-aware, and aligned with how a human proposal writer or business professional would communicate.

Follow these rules:

MAX COVER LETTER WORDS : 250 WORDS OR 1500 CHARACTER (whatever is minimum)

1. Start with:
   Kind Attention –
   Mr. {{executives}}
   M/s {{client_name}}

   Offer Ref: {{offer_ref}}; Date: {{letter_date}}

   Subject – Techno-Commercial Offer for {{project_title}}  

2. If there is only one executive, address them with:
   Dear {{first_exec_name}},
   If multiple executives, skip "Dear" and go directly to the content.
   Use Mr. for male and Ms. for female executives.

3. Opening paragraph:(Human way)
   - Acknowledge the invitation or requirement.
   - If invitation_date exists, mention it naturally.
   - If meeting_date exists, reference recent discussions or suggestions.
   - Wording must change between runs (not fixed sentences).
4. Body:(Human way)
   - Highlight Falcon’s analysis, solution evaluation, and technical proposal attachment.
   - Mention Falcon’s proven intralogistics technologies and experience.
   - Personalize with client_name.
   - Optionally mention project planning or timeline.
5. Closing: (Human way)
   - Reaffirm sender’s personal commitment.
   - Encourage the client to reach out for clarifications.
   - End with Best Regards, sender_name (sender_title).

Important:
- Do not exceed word/character limit.
- Keep tone formal, professional, and client-oriented.
- Do not copy exact sentences; rephrase wording across generations.
- The cover letter MUST be sounded human, natural and professional. It should be clear, authentic, and warm, without feeling robotic or overly formal.


Use the following examples as reference for tone and structure, DO NOT COPY EXACT EXAMPLES,rephrase example wording across generations:

Example 1:
Kind Attention –
Mr. Laurent Iem
Offer Ref: F25-00243-00;
Date: 07-08-2025
Subject – Techno-Commercial Offer for Cross Belt Sorter project 
Dear Laurent,
Thank you very much for your invitation of July 11th, to offer for the **sorter project in France and 
Spain**. Based on the meeting of July 17th and the suggestions provided by your technical, operational
and project teams over the past weeks, we have developed the detailed design for the sorter 
systems. Please find herewith our offer.
As you will note, we have done an in-depth data analysis and evaluated various solution options best 
suited for your requirements. Along with the information received, we have put together a **detailed 
technical proposal** laid out in various sections and sequenced to enable you to get a full insight in 
our proposed solution and to re-enforce our commitment to being your partner in this strategic 
initiative.
The core of the sortation system includes the well-known **Falcon horizontal electrical cross belt**
(dual belt on one carrier). 
In subsequent sections, we have highlighted the capabilities and experiences of Falcon Autotech 
with sections on our Intra-logistics Automation Technologies and references. 
Vinted has a global partner with local support. 
To conclude, I would like to add my personal commitment on behalf of Falcon Autotech. As we move 
through the RFP process, please do not hesitate to contact me and my team. We will be pleased to 
assist you with any further information or clarifications that you might have and look forward to our 
management call to finalise.
Best Regards on behalf of the team,
Johan Hoelen 
General Manager Europe


Example 2:
Kind Attention – 
Mr. Chetan Kumar 
Mr. Ravinder Sharma 
Mr. Vineet Kumar Pant 
Mr. Nikhil Devadas PV 
M/s Burjeel Holdings 
Offer Ref: F25-00201; Date: 14-06-2025 
Subject – Detailed Techno -Commercial Offer for Fulfilment Centre Automation for Burjeel 
Holdings 
We are pleased to submit our Techno-Commercial Offer in response to your requirement for 
Fulfilment centre automation for UAE, Abu Dhabi facility. 
As you will note, we have done an in-depth data analysis and evaluated various solution options best 
suited for your requirements, along with the information collected during the meetings and 
discussions with you, we have put together a detailed technical proposal laid out in various sections 
and sequenced to enable you to understand our proposed solution and to re-enforce our commitment 
to being your partner in this strategic initiative. 
In subsequent sections, we have highlighted the capabilities and experiences of Falcon Autotech with 
sections on our Intra-logistics Automation Technologies and references. 
To conclude, I would like to add my personal commitment on behalf of Falcon Autotech. As we move 
through the RFP process, please do not hesitate to contact me and my team. We will be pleased to 
assist you with any further information or clarifications that you might have. 
Best Regards, 
Sandeep Bansal 
Chief Business Officer


DO NOT ADD ANY EXTRA WORD OR INFO APART FROM THE COVER LETTER.
Highlight the main system or project name in main body (not subject line) as bold style, use **for Bold**.
""" 
PROMPT_EXEC_SUMMARY = """
You are a Proposal Writing Assistant specialized in Falcon Autotech automation projects.  
Falcon Autotech designs, manufactures, supplies, implements, and maintains warehouse automation solutions—such as sortation systems, conveyor automation, pick/put-to-light, ASRS robotics, and dimension & weight scanning—for industries including e-commerce, fashion, FMCG, pharma, groceries, and CE-P.  
The writing style must be indistinguishable from natural human writing. The text should read as if drafted by an experienced professional, not an AI system. Use clear, simple, and natural language with varied sentence lengths and structures. Avoid generic phrases, repetitive patterns, or mechanical tone. Ensure that the output flows smoothly, conveys intent naturally, and would not be detected as machine-generated. The content should feel thoughtful, context-aware, and aligned with how a human proposal writer or business professional would communicate.

Your task is to generate **unique, client-tailored Executive Summaries** based on the “Proposed System Description” section of Falcon proposals.  
The summary must always reflect Falcon’s style but **no two summaries should ever be identical**. Introduce subtle variations in wording, phrasing, and sentence structure while keeping the same professional tone.  

### Writing Rules

**Opening Section**
- Begin with Falcon Autotech’s commitment and strong interest in responding to the client’s requirement.  
- Mention Falcon’s partnership approach, customization, and proven track record.  
- Use varied sentence structures and synonyms so every generation feels different.  

**Bullet Points**
- Provide exactly **4–5 high-level system features or modules**.  
- Each bullet MUST be short, clear, and client-friendly (e.g., “Spiral Conveyors for smooth material flow”).  
- Avoid technical specifications, sub-bullets, or repeating the same idea in different words.  
- The order of bullets should vary slightly between generations.  
- Add numeric along with the components ONLY IF extensively mentioned in Proposed System Description
- Bold the main components of the system. There can be max 2-3 bold words.

**Closing Section**
- End with a **personalized closing statement**.  
- Reaffirm that the solution is tailored to meet the client’s technical and operational requirements.  
- Mention the RFP/customization and highlight benefits like efficiency, smooth material flow, and faster TAT.  
- Closing phrasing should change between runs (use variations in tone, sentence structure, and emphasis).  

### Important Constraints
- Keep the tone formal, professional, and benefit-driven.  
- Do **not** reuse exact sentences from earlier examples.  
- Ensure variability: two runs for the same input must never produce identical text.  
- Do **not** add any extra sections outside the defined structure.  

### Output Format
1. Opening paragraph (commitment + partnership).  
2. 4–6 bullet points (system modules).  
3. Closing personalized statement.  


**DO NOT ADD ANY EXTRA TEXT OR INFORMATION OR JUSTIFICATION or "Here is an Executive Summary for the proposal:" EXCEPT THE FULL PROPOSAL**

**Eaxmple 1:**

Falcon Autotech is pleased to confirm its great interest in responding to this Dnata RFQ of Conveyor Automation for 
Dubai Location. Our team has been working closely with the relevant stakeholders, with a clear commitment to 
listening, understanding your needs, and ensuring this project's success.
Following the same objective for the Dubai CBS system, we are happy to offer a compliant solution meeting all 
technical and operational requirements at competitive price, delivering key results.
Our solution is based on the following key characteristics:
    • Conveyor Automation for Handling Shipments of Boxes 
    • 3D ASRS (NEO) System for Storage of the shipments
    • Spiral Conveyors for smooth material flow
    • Outbound sorter with Swivel Wheel Divert Units for faster TAT and TPH
    • ULD Handling System (Castor Deck)
    • Oversize Cargo Pallet Conveyor Automation

This solution has been crafted especially keeping Dnata’s technical and operational requirements, as listed in the RFP 
document, making it a tailor-made solution delivering a faster TAT and an efficient material flow.


**Example 2 :**
Falcon is pleased to confirm its great interest in responding to this Mark3 International
Requirement of SWEDI Sorter. Our team has been working closely with the relevant 
stakeholders, with a clear commitment to listening and understanding your needs and 
ensuring this project's success.
Following the same objective for the Iran system, we are happy to offer a compliant solution 
meeting all technical and operational requirements, high-performance, optimized, tailormade, fast and secure planning, and a competitive price.
Our solution is based on the following key characteristics:

    • The SWEDI sorter, based on a Swivel wheel technology, offers a designed 
    throughput of 2,400 pph.
    • A Loading conveyor system with provision of manual Shipment loading.
    • 10 Output Chute & 1 Rejection/Sort Fail Chute

A tailor-made and simple layout designed explicitly for Mark3 International. The proposed 
layout is the result of the technical requirements in the RFP document and our discussions 
with the relevant stakeholders during our meeting.

**Example 3:**
Falcon is pleased to confirm its great interest in responding to this RFQ. Our team has been 
working closely with the relevant stakeholders, with a clear commitment to listening and 
understanding your needs and ensuring this project's success.
As prime contractor, Falcon ensures its full commitment to successfully completing this 
project. Following the same objective for the system, we are happy to offer a compliant 
solution meeting all technical and operational requirements, high-performance, optimized, 
tailor-made, fast and secure planning, and a competitive price. For full transparency, a 
compliancy list forms part of our offer. 
Our solution is based on the following key characteristics:

    • One Sorter, based on Falcon’s own, well-known loop cross-belt technology
    • 7 Sets of Induct with Manual Loading 
    • 1 Sets of Infeed with Volume Distribution System (VDS)
    • 40 Pcs Pallet Chute
    • 26 Pcs PTL Chutes for 312 Bags
    • 312 Pcs Put to Light with Rack for Bags
    • 3 Pcs Rejection Chute 
    • 7 Pcs of VDS Chute 
    • 7 Pcs of NC Chute

A tailor-made and simple layout, specifically designed to VINTED The proposed layout, is 
the result of the technical requirements in the RFP document and our discussions with the 
relevant stakeholders during the feedback round.
- Simple operational conditions due to one double decker horizontal cross belt sorter.
- Easy maintenance: optimized number of conveyors and concentrated inducts area.


DO NOT ADD ANY EXTRA WORD OR INFO APART FROM THE EXECUTIVE SUMMARY
"""
MERMAID_PROMPT = """
You are an expert systems analyst and Mermaid generator.
From SOLUTION_TEXT, extract the real process steps and output a vertical Mermaid flowchart.

RETURN FORMAT:
- Return ONLY Mermaid code that starts with: flowchart TD
- No markdown fences. No prose. No HTML tags in labels.
- Use short, unique IDs for all nodes (e.g., A, B1, C_out).
- Use standard shapes by best practice (rectangle for process, diamond for decision, oval for start/end, parallelogram for data/external).
- Keep layout vertical (top→down): TD.


**CLASS ASSIGNMENTS (required for every node):**
- main    = primary components (Infeed, VDS, Induct, Sorter/CBS/SWEDI, Scanning & Dimensioning, Print & Apply, etc.)
- sub     = secondary/auxiliary steps (staging, conveyors, data stores, handoffs, etc.)
- accept  = positive/approved/success/output/dispatch
- reject  = negative/fail/error/rejection/return-to-sender/technical

**STYLE BLOCK (append at the end EXACTLY as below) and blue arrows:**

classDef main fill:#5178B7,stroke:#5178B7,color:#ffffff,stroke-width:2px;
classDef sub fill:#FFC000,stroke:#B8860B,color:#111111,stroke-width:1.6px;
classDef accept fill:#96D157,stroke:#5A9831,color:#111111,stroke-width:2px;
classDef reject fill:#9C1922,stroke:#7F1420,color:#ffffff,stroke-width:2px;
linkStyle default stroke:#060c71,stroke-width:1.6px;

Then include one `class <ID> <classname>;` line for EVERY node ID you used (e.g., `class A main;`).

**CONSTRAINTS:**
- One main path top→bottom; branch only when the text clearly branches (e.g., accept vs reject).
- Normalize obvious synonyms (e.g., “PTL” → “Put to Light”), but keep the client’s nouns.
- No invented or duplicate nodes. No cycles unless explicitly described.
- Prefer concise labels; use \n for line breaks if needed (do not use HTML).
- Ensure the final output is parseable by Mermaid v10+.
- DO NOT decompose into micro-steps (belt segments, sensor ticks, minor handoffs) unless the text explicitly describes those sub-steps and they materially change sequence or routing.
- Every node label must be max 2–3 words, terse and non-descriptive. Prefer standard, unambiguous short forms when available (e.g., VDS, ICR, P&A, CBS, SWEDI, PTL). Avoid parentheses and filler words or two similar words. If no clear short form exists in the source, use the shortest clear wording.
- MUST ALWAYS use 'Put to Light' ONLY instead of PTL or Put to Light (PTL).
"""       

# ============================================================================
# Utilities
# ============================================================================
BULLET_RX = re.compile(r"^\s*([\u2022\-\*])\s+(.*)$")  # •  -, *
def _safe_bytes(file) -> Optional[bytes]:
    if not file:
        return None
    try:
        return file.getvalue()
    except Exception:
        return None

def _safe_name(file) -> Optional[str]:
    try:
        return file.name if file else None
    except Exception:
        return None

def _to_date_str(d) -> Optional[str]:
    try:
        return d.strftime("%d-%m-%Y") if d else None
    except Exception:
        return None

def _add_page_break(doc: Document) -> None:
    p = doc.add_paragraph()
    r = p.add_run()
    r.add_break(WD_BREAK.PAGE)

# ============================================================================
# PDF Parsing
# ============================================================================

def pdf_to_images_and_text(pdf_path: str) -> List[Dict[str, str]]:
    pages: List[Dict[str, str]] = []
    with fitz.open(pdf_path) as doc:
        for page in doc:
            text = page.get_text("text")
            pix = page.get_pixmap(dpi=200)
            img_bytes = pix.tobytes("png")
            img_b64 = base64.b64encode(img_bytes).decode("utf-8")
            pages.append({"text": text, "image_b64": img_b64})
    return pages

# ============================================================================
# Word Helpers (styling, headers, sections)
# ============================================================================

def add_response_heading(doc: Document, project_title: str) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f" {project_title}")
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0, 0, 0)

    highlight = OxmlElement("w:highlight")
    highlight.set(qn("w:val"), "lightGray")
    rPr = run._element.get_or_add_rPr()
    rPr.append(highlight)

def add_section_heading(doc: Document, title_text: str) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(title_text)
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.underline = True
    run.font.color.rgb = RGBColor(2, 12, 115)

def add_body_text(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.font.size = Pt(11)
    run.font.name = "Calibri"

def add_contact_box(doc: Document) -> Document:
    table = doc.add_table(rows=1, cols=1, style="Table Grid")
    cell = table.cell(0, 0)

    contact_text = (
        "Falcon Autotech Private Limited\n"
        "Plot No. 87, Sector Ecotech-1, Extension-1, Greater Noida, Uttar Pradesh 201308.\n\n"
        "Contact – Vaibhav Shukla\n"
        "Assistant Engineer\n"
        "Mob - +91 9871371775\n"
        "Vaibhav.shukla@falconautotech.com"
    )
    para = cell.paragraphs[0]
    run = para.add_run(contact_text)
    run.font.size = Pt(11)
    run.font.name = "Calibri"
    return doc

def add_hyperlink(paragraph, text: str, url: str, color: str = "000099", underline: bool = True):
    try:
        part = paragraph.part
        r_id = part.relate_to(
            url,
            reltype="http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True,
        )

        hyperlink = OxmlElement("w:hyperlink")
        hyperlink.set(qn("r:id"), r_id)

        new_run = OxmlElement("w:r")
        rPr = OxmlElement("w:rPr")

        if color:
            color_elem = OxmlElement("w:color")
            color_elem.set(qn("w:val"), color)
            rPr.append(color_elem)

        if underline:
            u_elem = OxmlElement("w:u")
            u_elem.set(qn("w:val"), "single")
            rPr.append(u_elem)

        rFonts = OxmlElement("w:rFonts")
        rFonts.set(qn("w:ascii"), "Calibri")
        rFonts.set(qn("w:hAnsi"), "Calibri")
        rFonts.set(qn("w:eastAsia"), "Calibri")
        rPr.append(rFonts)

        sz = OxmlElement("w:sz")
        sz.set(qn("w:val"), "16")  # 8pt
        rPr.append(sz)

        szCs = OxmlElement("w:szCs")
        szCs.set(qn("w:val"), "16")
        rPr.append(szCs)

        t = OxmlElement("w:t")
        t.text = text

        new_run.append(rPr)
        new_run.append(t)
        hyperlink.append(new_run)
        paragraph._element.append(hyperlink)
        return hyperlink
    except Exception:
        run = paragraph.add_run(f" {text} ")
        run.font.color.rgb = RGBColor(0, 0, 153)
        run.underline = True
        run.font.name = "Calibri"
        run.font.size = Pt(8)
        return None

def add_page_headers(
    doc: Document,
    client_logo: Optional[bytes],
    falcon_logo: Optional[bytes],
    client_name: str,
    project_name: str,
    make_first_page_different: bool = False,
) -> None:
    footer_text = "© FALCON AUTOTECH 2025 Confidential: Not for Distribution. "
    footer_url = "https://www.falconautotech.com/"

    if make_first_page_different and doc.sections:
        doc.sections[0].different_first_page_header_footer = True

    for section in doc.sections:
        section.header.is_linked_to_previous = False
        section.footer.is_linked_to_previous = False

        header = section.header
        footer = section.footer

        # Clear header
        try:
            for tbl in list(header.tables):
                tbl._element.getparent().remove(tbl._element)
            for p in list(header.paragraphs):
                p._element.getparent().remove(p._element)
        except Exception:
            pass

        table = header.add_table(rows=1, cols=3, width=Inches(6.5))
        table.autofit = False
        row = table.rows[0]
        row.cells[0].width = Inches(1.2)
        row.cells[1].width = Inches(4.1)
        row.cells[2].width = Inches(1.2)
        for cell in row.cells:
            cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER

        # Left: Client logo
        if client_logo:
            try:
                para = row.cells[0].paragraphs[0]
                run = para.add_run()
                run.add_picture(io.BytesIO(client_logo), width=Inches(1.0))
                para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            except Exception:
                pass

        # Middle: text
        header_text = f"FALCON’s Proposal to {client_name} for the {project_name}"
        para_mid = row.cells[1].paragraphs[0]
        run_mid = para_mid.add_run(header_text)
        run_mid.bold = False
        run_mid.font.size = Pt(8)
        run_mid.font.name = "Calibri"
        try:
            run_mid._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        except Exception:
            pass
        para_mid.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Right: Falcon logo
        if falcon_logo:
            try:
                para = row.cells[2].paragraphs[0]
                run = para.add_run()
                run.add_picture(io.BytesIO(falcon_logo), width=Inches(1.0))
                para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            except Exception:
                pass

        # Clear footer
        try:
            for para in list(footer.paragraphs):
                p_element = para._element
                p_element.getparent().remove(p_element)
        except Exception:
            pass

        para_f = footer.add_paragraph()
        run_f = para_f.add_run(footer_text)
        run_f.font.name = "Calibri"
        run_f.font.size = Pt(8)
        try:
            run_f._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        except Exception:
            pass
        add_hyperlink(para_f, footer_url, footer_url, color="000099", underline=True)
        para_f.alignment = WD_ALIGN_PARAGRAPH.CENTER

def add_runs_from_text(
    paragraph,
    text: str,
    italic: bool = False,
    base_font: str = "Calibri",
    base_size_pt: int = 11,
) -> None:
    pattern = re.compile(r"\*\*(.+?)\*\*")
    last = 0
    for m in pattern.finditer(text):
        pre = text[last:m.start()]
        if pre:
            run = paragraph.add_run(pre)
            run.font.name = base_font
            run.font.size = Pt(base_size_pt)
            run.italic = italic

        run = paragraph.add_run(m.group(1))
        run.bold = True
        run.italic = italic
        run.font.name = base_font
        run.font.size = Pt(base_size_pt)
        last = m.end()

    rem = text[last:]
    if rem:
        run = paragraph.add_run(rem)
        run.font.name = base_font
        run.font.size = Pt(base_size_pt)
        run.italic = italic

def save_to_docx(
    content: str,
    title: str,
    client_logo: Optional[bytes],
    falcon_logo: Optional[bytes],
    client_name: str,
    project_name: str,
) -> io.BytesIO:
    doc = Document()
    add_page_headers(doc, client_logo, falcon_logo, client_name, project_name)

    if title.strip():
        title_para = doc.add_paragraph()
        run_title = title_para.add_run(title)
        run_title.font.name = "Calibri"
        run_title.font.size = Pt(14)
        run_title.bold = True
        run_title.underline = True
        run_title.font.color.rgb = RGBColor(51, 51, 153)
        title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT

    for para in content.split("\n"):
        if not para.strip():
            continue

        line = para.strip()

        if line.startswith('•') or line.startswith('·') or line.startswith('*'):
            bullet_para = doc.add_paragraph(style='List Bullet')
            bullet_text = line[1:].strip()
            add_runs_from_text(
                bullet_para,
                bullet_text,
                italic=False,
                base_font="Calibri",
                base_size_pt=11,
            )
            bullet_para.paragraph_format.left_indent = Inches(0.5)
            bullet_para.paragraph_format.first_line_indent = Inches(-0.5)

        else:
            p = doc.add_paragraph()
            special = (
                line.startswith("Kind Attention")
                or line.startswith("Offer Ref")
                or line.startswith("Date")
                or line.startswith("Mr.")
                or line.startswith("MR.")
                or line.startswith("Ms.")
                or line.startswith("MS.")
                or line.startswith("M/s")
                or line.startswith("M/S")
                or line.startswith("Chief")
                or line.startswith("Sandeep")
                or line.startswith("Subject")
                or line.startswith("Best Regards")
            )
            if special:
                run = p.add_run(line)
                run.font.name = "Calibri"
                try:
                    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
                except Exception:
                    pass
                run.font.size = Pt(11)
                run.bold = True
            else:
                add_runs_from_text(p, line, italic=False, base_font="Calibri", base_size_pt=11)

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def create_cover_page(
    template_path: str,
    client_logo: Optional[bytes],
    client_name: str,
    project_title: str,
) -> io.BytesIO:
    doc = Document(template_path)

    for sec in doc.sections:
        for part in (
            getattr(sec, "header", None),
            getattr(sec, "footer", None),
            getattr(sec, "first_page_header", None),
            getattr(sec, "first_page_footer", None),
            getattr(sec, "even_page_header", None),
            getattr(sec, "even_page_footer", None),
        ):
            if not part:
                continue
            try:
                part.is_linked_to_previous = False
            except Exception:
                pass
            try:
                for tbl in list(part.tables):
                    tbl._element.getparent().remove(tbl._element)
                for p in list(part.paragraphs):
                    p._element.getparent().remove(p._element)
            except Exception:
                pass

    if client_logo:
        try:
            im = Image.open(io.BytesIO(client_logo))
            if im.mode != "RGBA":
                im = im.convert("RGBA")
            alpha = im.getchannel("A")
            bbox = alpha.getbbox()
            if bbox:
                im = im.crop(bbox)
                alpha = im.getchannel("A")
            bg = Image.new("RGB", im.size, (255, 255, 255))
            bg.paste(im, mask=alpha)

            buf = io.BytesIO()
            bg.save(buf, format="PNG")
            buf.seek(0)

            first_para = doc.paragraphs[0]
            run_logo = first_para.insert_paragraph_before().add_run()
            run_logo.add_picture(buf, width=Inches(2.0))
        except Exception:
            first_para = doc.paragraphs[0]
            run_logo = first_para.insert_paragraph_before().add_run()
            run_logo.add_picture(io.BytesIO(client_logo), width=Inches(2.0))

    for _ in range(6):
        doc.add_paragraph("")

    title = f"FALCON’s Proposal to {client_name} for the {project_title}"
    p = doc.add_paragraph()
    run = p.add_run(title)
    run.font.size = Pt(24)
    run.font.bold = False
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT

    today_str = datetime.today().strftime("%B %d, %Y")
    p2 = doc.add_paragraph()
    run2 = p2.add_run(today_str)
    run2.font.size = Pt(14)
    run2.font.name = "Calibri"
    run2.font.color.rgb = RGBColor(255, 215, 0)
    p2.alignment = WD_ALIGN_PARAGRAPH.LEFT

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def create_rfq_response_doc(
    project_title: str, ref_no: str, image_bytes: Optional[bytes] = None
) -> io.BytesIO:
    doc = Document()
    add_response_heading(doc, f"Response to Requirement for {project_title}")
    doc.add_paragraph("")
    if image_bytes:
        try:
            doc.add_picture(io.BytesIO(image_bytes), width=Inches(6))
        except Exception:
            pass
        doc.add_paragraph("")
    add_response_heading(doc, f"Proposal Reference – {ref_no}")
    doc.add_paragraph("")
    add_contact_box(doc)

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def merge_docx_files_with_page_breaks(buffers_list: List[Optional[io.BytesIO]]) -> io.BytesIO:
    if not buffers_list or buffers_list[0] is None:
        raise ValueError("At least one valid document buffer is required as base document.")

    base_doc = Document(buffers_list[0])
    composer = Composer(base_doc)

    for buf in buffers_list[1:]:
        if buf is None:
            continue
        _add_page_break(composer.doc)
        doc_to_append = Document(buf)
        composer.append(doc_to_append)

    merged_buffer = io.BytesIO()
    composer.save(merged_buffer)
    merged_buffer.seek(0)
    return merged_buffer

# ============================================================================
# Company Profile Builder (unchanged content)
# ============================================================================
def build_company_profile(doc: Document) -> None:
    """
    Populate the given `doc` with the original multi-page Company Profile content.
    Mirrors the structure, text, image sizes, and page breaks from your source.
    """

    # ---------- Page 1 ----------
    add_section_heading(doc, "3. Company Profile")

    top_text = (
    "Falcon Autotech (Falcon) is a global intralogistics automation solutions company. "
    "With over 10 years of experience, Falcon has worked with some of the most innovative "
    "brands in E-Commerce, CEP, Fashion, Food/FMCG, Auto and Pharmaceutical Industries. "
    "With our proprietary software and robust hardware integration capabilities, Falcon designs, "
    "manufactures, supplies, implements, and maintains world-class warehouse automation systems globally. "
    "Falcon’s strong research and development team and the continuous focus on innovation reflect our strong "
    "solution line around Sortation, Robotics, Conveying, Vision Systems and IOT. "
    "Falcon has done over 1,800 installations across 15 countries on four continents."
)
    add_body_text(doc, top_text)

    # Insert Image (portrait, ~6 inches)
    img_path1 = "Input\\Static_AboutCompany\\1.png"  # <-- Replace with your actual image file path
    if os.path.exists(img_path1):
        p = doc.add_paragraph()
        run = p.add_run()
        run.add_picture(img_path1, width=Inches(6))
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    

    bottom_text = (
    "Falcon Autotech is currently among the top 15 intralogistics automation company, our vision is to become top 10 intralogistics automation company in our focused product lines."
        
    )
    add_body_text(doc, bottom_text)
    # Insert Image (portrait, ~6 inches)
    img_path2 = "Input\\Static_AboutCompany\\2.png"  # <-- Replace with your actual image file path
    if os.path.exists(img_path2):
        p = doc.add_paragraph()
        run = p.add_run()
        run.add_picture(img_path2, width=Inches(5))
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
    doc.add_page_break()

    # ---------- Page 2 ----------
    
    top_text2 = (
    "The team started out in 2004 solving special purpose automation problems for clients and later established Falcon Autotech in 2012 with strong focus on building standard technology stack spanning across Hardware, Firmware and Software to tackle bigger Supply Chain problems around warehouse automation and material handling."
    "Over the decade, Falcon has made rapid strides and has carved out a niche in some of the world's most cutting-edge technologies: Sortation, Robotics, Conveying, Vision Systems and IOT."
)
    add_body_text(doc, top_text2)

    # Insert Image (portrait, ~6 inches)
    img_path3 = "Input\\Static_AboutCompany\\3.png"  # <-- Replace with your actual image file path
    if os.path.exists(img_path3):
        p = doc.add_paragraph()
        run = p.add_run()
        run.add_picture(img_path3, width=Inches(6))
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    

    bottom_text2 = (
    "As a leading player in the intra-logistics automation space, Falcon continuously strives to improve the operational efficiencies and accuracies for its clients through its domain knowledge and experience in addition to its wide range of products and solutions."
    "In order to be able to live up to the high expectations set forth by our clients, the team at Falcon realizes the importance of taking up selective applications in focused Industries and deliver world class projects in return."
    )
    add_body_text(doc, bottom_text2)
    # Insert Image (portrait, ~6 inches)
    img_path4 = "Input\\Static_AboutCompany\\4.png"  # <-- Replace with your actual image file path
    if os.path.exists(img_path4):
        p = doc.add_paragraph()
        run = p.add_run()
        run.add_picture(img_path4, width=Inches(6))
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        

   

    # ---------- Page 3 ----------
    
    img_path5 = "Input\\Static_AboutCompany\\5.png"  # <-- Replace with your actual image file path
    if os.path.exists(img_path5):
        p = doc.add_paragraph()
        run = p.add_run()
        run.add_picture(img_path5, width=Inches(6))
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    bottom_text3 = (
    "Falcon Autotech has successfully delivered warehouse automation solutions based on smart and innovative combinations of above product lines for effective materials handling, sortation and movement."
    "The process is controlled in real-time by our in house WCS applications."
    "These solutions considerably cut the need for manual operations, improve working conditions and ensure the highest accuracy of the entire process up to final delivery to the recipient."
    "\n"
    "Over the last 10 years, Falcon has worked with some of the most innovative brands worldwide and has established long standing partnerships."
    " These brands are testimony of our strong focus on delivering superior customer satisfaction and offering end-to-end intralogistics solutions."
    )
    add_body_text(doc, bottom_text3)

    img_path6 = "Input\\Static_AboutCompany\\6.png"  # <-- Replace with your actual image file path
    if os.path.exists(img_path6):
        p = doc.add_paragraph()
        run = p.add_run()
        run.add_picture(img_path6, width=Inches(6))
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_page_break()

    # ---------- Page 4 ----------
    bottom_text4 = (
    "With over 1,800 installations, today Falcon’s systems are used all over the globe. Falcon has highly motivated team of 600+ employees supported by over 15 global partners who help us design, manufacture, deliver and maintain automation solutions globally."
    )

    add_body_text(doc, bottom_text4)
    img_path7 = "Input\\Static_AboutCompany\\7.png"  # <-- Replace with your actual image file path
    if os.path.exists(img_path7):
        p = doc.add_paragraph()
        run = p.add_run()
        run.add_picture(img_path7, width=Inches(6))
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_section_heading(doc, "Customer Engagement Model")
    
    img_path8 = "Input\\Static_AboutCompany\\8.png"  # <-- Replace with your actual image file path
    if os.path.exists(img_path8):
        p = doc.add_paragraph()
        run = p.add_run()
        run.add_picture(img_path8, width=Inches(7))
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_page_break()

    # ---------- Page 5 ----------
    add_section_heading(doc, "3. Falcon’s Experience and Achievements in Sortation Space Globally")
    bullet_points = [
    "Ranked among Top 20 Sortation System Suppliers globally.",
    "Currently possess one of the World’s largest portfolios in Sortation Technologies: 7 In-house technologies.",
    "Total installed capacity of 10 million Shipments per day worldwide.",
    "Only company to be able to offer a Fully Integrated AMS."
]

    for point in bullet_points:
        para = doc.add_paragraph(point, style="List Bullet")
        run = para.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(10)

    img_path9 = "Input\\Static_AboutCompany\\9.png"  # <-- Replace with your actual image file path
    if os.path.exists(img_path9):
        p = doc.add_paragraph()
        run = p.add_run()
        run.add_picture(img_path9, width=Inches(6))
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    

    img_path10 = "Input\\Static_AboutCompany\\10.png"  # <-- Replace with your actual image file path
    if os.path.exists(img_path10):
        p = doc.add_paragraph()
        run = p.add_run()
        run.add_picture(img_path10, width=Inches(6))
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
# ---------- Page 5 ----------
    img_path11 = "Input\\Static_AboutCompany\\11.png"  # <-- Replace with your actual image file path
    if os.path.exists(img_path11):
        p = doc.add_paragraph()
        run = p.add_run()
        run.add_picture(img_path11, width=Inches(6))
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("")
    img_path12 = "Input\\Static_AboutCompany\\12.png"  # <-- Replace with your actual image file path
    if os.path.exists(img_path12):
        p = doc.add_paragraph()
        run = p.add_run()
        run.add_picture(img_path12, width=Inches(6))
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
# ============================================================================
# LLM Generators
# ============================================================================

def _groq_cover_letter_call(
    system_prompt: str,
    user_prompt: str,
    model: str = "openai/gpt-oss-120b",
    temperature: float = 0.7,
    top_p: float = 0.9,
    max_tokens: int = 600,
) -> str:
    if groq_client is None:
        raise RuntimeError("Groq client is not initialized. Check GROQ_API_KEY.")
    resp = groq_client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=temperature,
        top_p=top_p,
        max_tokens=max_tokens,
    )
    return resp.choices[0].message.content

def _groq_exec_summary_call(
    system_prompt: str,
    user_prompt: str,
    pdf_text: str,
    model: str = "meta-llama/llama-4-scout-17b-16e-instruct",
    temperature: float = 0.7,
    top_p: float = 0.9,
    max_tokens: int = 1000,
) -> str:
    if groq_client is None:
        raise RuntimeError("Groq client is not initialized. Check GROQ_API_KEY.")
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt + pdf_text},
    ]
    resp = groq_client.chat.completions.create(
        model=model,
        messages=messages,
        temperature=temperature,
        top_p=top_p,
        max_tokens=max_tokens,
    )
    return resp.choices[0].message.content

def generate_cover_letters_concurrent(
    sender_name: str,
    sender_title: str,
    client_name: str,
    executives: str,
    offer_ref: str,
    project_title: str,
    invitation_date: Optional[str],
    meeting_date: Optional[str],
) -> Tuple[str, str]:
    system_prompt = PROMPT_COVER_LETTER
    letter_date = datetime.now().strftime("%d-%m-%Y")
    user_prompt = f"""
Generate a formal techno-commercial COVER LETTER for a proposal from below INPUT.

Client: {client_name}
Executives: {executives}
Offer Ref: {offer_ref}
Date: {letter_date}
Project: {project_title}
Invitation Date: {invitation_date}
Meeting Date: {meeting_date}
Sender Name: {sender_name}
Sender Title: {sender_title}
"""
    outputs: List[Optional[str]] = [None, None]
    with ThreadPoolExecutor(max_workers=2) as ex:
        futures = {
            ex.submit(_groq_cover_letter_call, system_prompt, user_prompt): 0,
            ex.submit(_groq_cover_letter_call, system_prompt, user_prompt): 1,
        }
        for fut in as_completed(futures):
            idx = futures[fut]
            try:
                outputs[idx] = fut.result()
            except Exception as e:
                outputs[idx] = f"Cover letter generation failed: {e}"
    return outputs[0] or "", outputs[1] or ""

def generate_cover_letter(
    sender_name: str,
    sender_title: str,
    client_name: str,
    executives: str,
    offer_ref: str,
    project_title: str,
    invitation_date: Optional[str] = None,
    meeting_date: Optional[str] = None,
    letter_date: Optional[str] = None,
) -> str:
    system_prompt = PROMPT_COVER_LETTER
    letter_date = letter_date or datetime.now().strftime("%d-%m-%Y")
    user_prompt = f"""
Generate a formal techno-commercial COVER LETTER for a proposal.

Client: {client_name}
Executives: {executives}
Offer Ref: {offer_ref}
Date: {letter_date}
Project: {project_title}
Invitation Date: {invitation_date}
Meeting Date: {meeting_date}
Sender Name: {sender_name}
Sender Title: {sender_title}
"""
    try:
        return _groq_cover_letter_call(system_prompt, user_prompt)
    except Exception as e:
        return f"Cover letter generation failed: {e}"

def generate_exec_summaries_concurrent(
    pdf_path: str, client_name: str, project_name: str
) -> Tuple[str, str]:
    system_prompt = PROMPT_EXEC_SUMMARY
    user_prompt = f"""
Generate an Executive Summary for the following proposal.

Client: {client_name}
Project: {project_name}
"""
    try:
        pages = pdf_to_images_and_text(pdf_path)
        if not pages:
            return "No pages found in PDF.", "No pages found in PDF."
        pdf_text = "\n".join(p["text"] for p in pages if p["text"].strip())
    except Exception as e:
        return f"PDF parsing failed: {e}", f"PDF parsing failed: {e}"

    outputs: List[Optional[str]] = [None, None]
    with ThreadPoolExecutor(max_workers=2) as ex:
        futures = {
            ex.submit(_groq_exec_summary_call, system_prompt, user_prompt, pdf_text): 0,
            ex.submit(_groq_exec_summary_call, system_prompt, user_prompt, pdf_text): 1,
        }
        for fut in as_completed(futures):
            idx = futures[fut]
            try:
                outputs[idx] = fut.result()
            except Exception as e:
                outputs[idx] = f"Executive summary generation failed: {e}"
    return outputs[0] or "", outputs[1] or ""

def generate_executive_summary(pdf_path: str, client_name: str, project_name: str) -> str:
    system_prompt = PROMPT_EXEC_SUMMARY
    user_prompt = f"""
Generate an Executive Summary for the following proposal.

Client: {client_name}
Project: {project_name}
"""
    try:
        pages = pdf_to_images_and_text(pdf_path)
        if not pages:
            return "No pages found in PDF."
        pdf_text = "\n".join(p["text"] for p in pages if p["text"].strip())
    except Exception as e:
        return f"PDF parsing failed: {e}"
    try:
        return _groq_exec_summary_call(system_prompt, user_prompt, pdf_text)
    except Exception as e:
        return f"Executive summary generation failed: {e}"

# ============================================================================
# FLOWCHART (Mermaid) — PDF→Text, LLM → Mermaid, Render, PPTX, DOCX
# ============================================================================

SECTION_HINTS = [
    "Proposed System Description",
    "Process flow of the System",
    "Process Flow of the System",
    "Summary of the System",
    "Concept Description",
    "Layout Overview",
    "Inbound", "Outbound", "Sorting", "NEO", "PTL", "CBS", "SWEDI", "Process flow:"
]

def _page_words_to_lines(page: fitz.Page, y_tol: float = 2.0) -> List[str]:
    words = page.get_text("words") or []
    if words and len(words[0]) >= 8:
        words.sort(key=lambda w: (w[5], w[6], w[0]))
        grouped = {}
        for x0, y0, x1, y1, txt, blk, ln, wn in words:
            key = (blk, ln)
            grouped.setdefault(key, []).append((x0, txt))
        lines = []
        for key in sorted(grouped.keys()):
            tokens = [t for _, t in sorted(grouped[key], key=lambda z: z[0])]
            lines.append(" ".join(tokens))
        return lines

    words.sort(key=lambda w: (w[1], w[0]))
    lines = []
    curr_y, curr = None, []
    for x0, y0, x1, y1, txt, *_ in words:
        if curr_y is None or abs(y0 - curr_y) <= y_tol:
            curr.append((x0, txt))
            curr_y = y0 if curr_y is None else (curr_y + y0) / 2.0
        else:
            lines.append(" ".join(t for _, t in sorted(curr, key=lambda z: z[0])))
            curr = [(x0, txt)]
            curr_y = y0
    if curr:
        lines.append(" ".join(t for _, t in sorted(curr, key=lambda z: z[0])))
    return lines

def _ocr_page_to_text(page: fitz.Page, dpi: int = 250) -> str:
    if not HAS_TESS:
        return ""
    pix = page.get_pixmap(dpi=dpi, alpha=False)
    img = PILImage.frombytes("RGB", [pix.width, pix.height], pix.samples)
    return pytesseract.image_to_string(img) or ""

def extract_relevant_text(pdf_bytes: bytes, max_pages: int = 15) -> str:
    text_chunks: List[str] = []
    with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
        pages = min(max_pages, doc.page_count)
        for i in range(pages):
            page = doc.load_page(i)
            lines = _page_words_to_lines(page)
            page_text = "\n".join(lines).strip()
            if len(page_text) < 40:
                ocr_txt = _ocr_page_to_text(page)
                if len(ocr_txt) > len(page_text):
                    page_text = ocr_txt
            text_chunks.append(page_text)

    blob = "\n".join(text_chunks)
    blob = re.sub(r"(\w)-\n(\w)", r"\1\2", blob)
    blob = re.sub(r"([A-Za-z])\s{1,2}([A-Za-z])", r"\1\2", blob)
    blob = re.sub(r"\s{2,}", " ", blob)
    blob = re.sub(r"[ \t]+\n", "\n", blob)
    blob = re.sub(r"\n{3,}", "\n\n", blob)

    blocks = []
    for hint in SECTION_HINTS:
        for m in re.finditer(rf"{re.escape(hint)}(.+?)(?:\n[A-Z][^\n]{{0,60}}\n|$)", blob, flags=re.IGNORECASE | re.DOTALL):
            seg = m.group(0)
            if len(seg) > 200:
                blocks.append(seg)
    chosen = "\n\n".join(blocks) if blocks else blob

    m_list = re.search(r"Process flow[:：-]?\s*(.+?)(?:\n\n|$)", chosen, flags=re.IGNORECASE | re.DOTALL)
    if m_list and len(m_list.group(1)) > 120:
        chosen = m_list.group(1)

    return chosen.strip()

def load_groq() -> Optional[Groq]:
    key = (os.getenv("GROQ_API_KEY") or "").strip()
    if not key:
        return None
    try:
        return Groq(api_key=key)
    except Exception:
        return None

def ask_llm_mermaid(client: Groq, solution_text: str) -> str:
    user = f"SOLUTION_TEXT:\n{solution_text}\n\nReturn only Mermaid code (start with: flowchart TD)."
    resp = client.chat.completions.create(
        model="meta-llama/llama-4-scout-17b-16e-instruct",
        messages=[{"role": "system", "content": MERMAID_PROMPT},
                  {"role": "user", "content": user}],
        temperature=0.1,
        max_tokens=1200,
    )
    out = (resp.choices[0].message.content or "").strip()
    out = re.sub(r"^```(?:mermaid)?\s*", "", out)
    out = re.sub(r"\s*```$", "", out)
    if not out.lower().startswith("flowchart"):
        m = re.search(r"(flowchart\s+TD[\s\S]+)$", out, re.IGNORECASE)
        if m:
            out = m.group(1).strip()
    return out

# ---- Mermaid rendering helpers ----
def sanitize_mermaid_for_render(code: str) -> str:
    if not code:
        return code
    s = code.replace("\r\n", "\n").replace("\r", "\n")

    def _fix(s: str, left: str, right: str) -> str:
        pat = re.compile(re.escape(left) + r"(.*?)" + re.escape(right), re.DOTALL)
        def sub(m):
            inner = m.group(1)
            inner = inner.replace("\\n", "<br/>").replace("\n", "<br/>")
            inner = re.sub(r"\s{2,}", " ", inner).strip()
            return f"{left}{inner}{right}"
        return pat.sub(sub, s)

    s = _fix(s, "[[", "]]")
    s = _fix(s, "((", "))")
    s = _fix(s, "[/", "/]")
    s = _fix(s, "[", "]")
    s = _fix(s, "(", ")")
    s = _fix(s, "{", "}")
    return s

def render_mermaid_chart(mermaid_code: str, height: int = 600):
    code_js = json.dumps(mermaid_code or "flowchart TD\n  A[Empty]")
    html = ''.join([
        '<div id="mmd_wrapper" style="width:100%; height:100%; overflow:hidden; '
        'border:1px solid #5178B7; border-radius:8px; background:#fff;">',
        '  <div id="mmd_container" style="width:100%; height:100%; position:relative;"></div>',
        '</div>',
        '<script src="https://cdn.jsdelivr.net/npm/mermaid@10/dist/mermaid.min.js"></script>',
        '<script>',
        'var code = ' + code_js + ';',
        'var container = document.getElementById("mmd_container");',
        'mermaid.initialize({ startOnLoad: false, securityLevel: "loose", ',
        '  flowchart: { useMaxWidth: true, htmlLabels: false, nodeSpacing: 20, rankSpacing: 30, padding: 8 },',
        '  theme: "default",',
        '  themeVariables: { fontSize: "12px" }',
        '});',
        'mermaid.render("theGraph", code).then(function(res){',
        '  container.innerHTML = res.svg;',
        '  var s = container.querySelector("svg");',
        '  if (s){',
        '    if (!s.getAttribute("viewBox")) {',
        '      try { var bb = s.getBBox(); s.setAttribute("viewBox", "0 0 " + Math.ceil(bb.width) + " " + Math.ceil(bb.height)); } catch(e) {}',
        '    }',
        '    s.removeAttribute("width");',
        '    s.removeAttribute("height");',
        '    s.setAttribute("preserveAspectRatio", "xMidYMid meet");',
        '    s.style.width = "100%";',
        '    s.style.height = "100%";',
        '    s.style.display = "block";',
        '  }',
        '}).catch(function(err){',
        '  container.innerHTML = "<pre style=\\"white-space:pre-wrap; color:#b00020;\\">" + ',
        '    (err && err.message ? err.message : String(err)) + "</pre>";',
        '});',
        '</script>',
    ])
    components.html(html, height=height, scrolling=True)

# ---- Mermaid→PNG (multi-strategy: Kroki → SVG→PNG → mermaid.ink → Selenium) ----
def kroki_url() -> str:
    url = (os.getenv("KROKI_URL") or "").strip()
    return url if url else "https://kroki.io"

def _deflate_b64_urlsafe(s: str) -> str:
    comp = zlib.compress(s.encode("utf-8"))
    return base64.urlsafe_b64encode(comp).decode("ascii")

def _svg_bytes_to_png(svg_bytes: bytes) -> Optional[bytes]:
    if not svg_bytes:
        return None
    if HAS_CAIROSVG:
        try:
            return cairosvg.svg2png(bytestring=svg_bytes)
        except Exception:
            pass
    return None

def mermaid_to_png_via_kroki(mermaid_code: str) -> Optional[bytes]:
    if not HAS_REQUESTS or not mermaid_code:
        return None
    try:
        url = kroki_url().rstrip("/") + "/mermaid/png"
        r = requests.post(url, data=mermaid_code.encode("utf-8"),
                          headers={"Content-Type": "text/plain"}, timeout=25)
        if r.ok and r.content:
            return r.content
    except Exception:
        pass

    try:
        url_svg = kroki_url().rstrip("/") + "/mermaid/svg"
        r = requests.post(url_svg, data=mermaid_code.encode("utf-8"),
                          headers={"Content-Type": "text/plain"}, timeout=25)
        if r.ok and r.content:
            png = _svg_bytes_to_png(r.content)
            if png:
                return png
    except Exception:
        pass

    try:
        code_b64 = _deflate_b64_urlsafe(mermaid_code)
        url_png = f"https://mermaid.ink/img/{code_b64}"
        r = requests.get(url_png, timeout=25, headers={"Accept": "image/png"})
        if r.ok and r.content:
            return r.content
    except Exception:
        pass

    try:
        code_b64 = _deflate_b64_urlsafe(mermaid_code)
        url_svg = f"https://mermaid.ink/svg/{code_b64}"
        r = requests.get(url_svg, timeout=25, headers={"Accept": "image/svg+xml"})
        if r.ok and r.content:
            png = _svg_bytes_to_png(r.content)
            if png:
                return png
    except Exception:
        pass

    return None

# ---- (NEW) Selenium-based fallback PNG renderer (from your Flowchart code) ----
def mermaid_to_png_via_chrome(mermaid_code: str) -> Optional[bytes]:
    if not (HAS_SELENIUM and mermaid_code):
        return None
    html_content = f"""<!doctype html>
<html><head><meta charset="utf-8">
<style>body{{margin:0;background:#fff;display:flex;align-items:center;justify-content:center;min-height:100vh}}
.wrap{{padding:20px}}#m.mermaid{{max-width:1600px}}</style>
</head><body>
<div class="wrap"><div class="mermaid" id="m">{mermaid_code}</div></div>
<script src="https://cdn.jsdelivr.net/npm/mermaid@10/dist/mermaid.min.js"></script>
<script>mermaid.initialize({{ startOnLoad: true, theme:"default" }});</script>
</body></html>"""
    try:
        import tempfile, time
        tmpdir = tempfile.mkdtemp(prefix="mmd2png_")
        html_path = os.path.join(tmpdir, "diagram.html")
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(html_content)

        options = webdriver.ChromeOptions()
        options.add_argument("--headless=new")
        options.add_argument("--disable-gpu")
        options.add_argument("--no-sandbox")
        options.add_argument("--disable-dev-shm-usage")
        options.add_argument("--start-maximized")

        driver = webdriver.Chrome(service=ChromeService(ChromeDriverManager().install()), options=options)
        driver.set_window_size(1920, 1200)
        driver.get(f"file://{html_path}")
        time.sleep(1.6)

        png = driver.get_screenshot_as_png()
        try:
            driver.quit()
        except Exception:
            pass
        return png
    except Exception:
        return None

def mermaid_to_png_best_effort(code: str) -> Optional[bytes]:
    png = mermaid_to_png_via_kroki(code)
    if png:
        return png
    png = mermaid_to_png_via_chrome(code)
    return png

# ---- Build a 1-page DOCX from PNG (or fallback to code) ----
def make_docx_fit_one_page_from_png_or_code(png_bytes: bytes,
                                            mermaid_code: str,
                                            title="Concept Description",
                                            client_logo=None, falcon_logo=None,
                                            client_name="", project_name="") -> io.BytesIO:
    png = _require_valid_png(png_bytes)
    return _docx_from_png_onepage(
        png, title,
        client_logo=client_logo, falcon_logo=falcon_logo,
        client_name=client_name, project_name=project_name
    )


# ---- Mermaid→PPTX (editable) ----
import collections

def _pptx_rgb(hexstr: str):
    hexstr = hexstr.lstrip("#")
    return PPTXRGB(int(hexstr[0:2],16), int(hexstr[2:4],16), int(hexstr[4:6],16))

def _parse_mermaid_nodes_edges(code: str):
    node_re = re.compile(r'^\s*([A-Za-z][\w]*)\s*(\(\(|\(|\[\[|\[|\{|\[\/)\s*([^)\]\}]+?)\s*(\)\)|\)|\]\]|\]|\}|\/\])', re.M)
    edge_re = re.compile(r'^\s*([A-Za-z][\w]*)\s*[-.]*>\s*(?:\|[^|]*\|\s*)?([A-Za-z][\w]*)', re.M)
    class_re = re.compile(r'^\s*class\s+([A-Za-z][\w]*)\s+([A-Za-z_][\w]*);', re.M)

    nodes = []
    for m in node_re.finditer(code):
        nid, l, label, r = m.groups()
        t = l + r
        nodes.append((nid, label.strip().replace("<br/>", "\n"), t))

    edges = [(m.group(1), m.group(2)) for m in edge_re.finditer(code)]
    classes = {m.group(1): m.group(2) for m in class_re.finditer(code)}
    return nodes, edges, classes

def _shape_type(token: str):
    token = (token or "").strip()
    if token in ("[//]", "[/]"):
        token = "[/]"
    return {
        "(())": SHAPE.FLOWCHART_TERMINATOR,
        "()":   SHAPE.FLOWCHART_PROCESS,
        "[[]]": SHAPE.FLOWCHART_PREDEFINED_PROCESS,
        "[]":   SHAPE.FLOWCHART_PROCESS,
        "{}":   SHAPE.FLOWCHART_DECISION,
        "[/]":  SHAPE.FLOWCHART_DATA,
    }.get(token, SHAPE.FLOWCHART_PROCESS)

def _infer_class(label: str):
    L = label.lower()
    if any(k in L for k in ("reject","rejection","fail","error","technical","rts","return to sender")):
        return "reject"
    if any(k in L for k in ("accept","approved","success","output","dispatch")):
        return "accept"
    if any(k in L for k in ("infeed","volume distribution","vds","induct","sorter","cbs","swedi","scanning","print & apply","dimension","weigh")):
        return "main"
    return "sub"



# Convenience: ensure we have flowchart DOCX buffer in session (to include in final merge)
# Replace your existing ensure_flowchart_docx_in_session with this DRAW.IO aware version
def ensure_flowchart_docx_in_session():
    """Build a 1-page DOCX from the current diagram (prefer latest Draw.io export)."""
    code = st.session_state.get("mermaid_code", "")
    _ensure_flow_state()
    fs = st.session_state.flow_state

    png_bytes = None
    # Prefer a fresh export if Draw.io is active
    if fs.mode == "drawio" and fs.driver is not None:
        try:
            png_bytes = _export_latest_flowchart_png()
        except Exception:
            png_bytes = None

    # Otherwise use the Kroki path
    if png_bytes is None and code.strip():
        png_bytes = mermaid_to_png_via_kroki(code)

    st.session_state["mermaid_png_bytes"] = png_bytes

    flow_docx = make_docx_fit_one_page_from_png_or_code(
        png_bytes,
        code,
        title="Concept Description",
        client_logo=st.session_state.get("client_logo"),
        falcon_logo=st.session_state.get("falcon_logo"),
        client_name=st.session_state.get("client_name",""),
        project_name=st.session_state.get("project_title",""),
    )
    st.session_state.flowchart_docx_buffer = flow_docx
    return png_bytes is not None


# ============================================================================
# Enhanced Streamlit UI - Falcon Proposal Automation (+ Flowchart Step)
# ============================================================================

st.set_page_config(
    page_title="Sales Proposal Automation System", 
    layout="wide",
    initial_sidebar_state="collapsed"
)

# Ensure Falcon logo bytes are available app-wide
if "falcon_logo" not in st.session_state:
    try:
        with open("Input\\Static_AboutCompany\\Falcon-Autotech_Logo.png", "rb") as _f:
            st.session_state.falcon_logo = _f.read()
    except Exception:
        st.session_state.falcon_logo = None

# Enhanced Professional Styling (Theme)
st.markdown("""
<style>
    .main > div { padding-top: 2rem; padding-bottom: 2rem; }
    .main-header {
        background: #060c71;
        background: linear-gradient(90deg, #060c71 0%, #2a3bb8 35%, #f9d20e 100%);
        padding: 2rem; border-radius: 15px; margin-bottom: 2rem;
        box-shadow: 0 8px 32px rgba(6, 12, 113, 0.3); color: white;
    }
    .main-header h1 { color: #fff !important; font-size: 2.5rem; font-weight: 700; margin: 0; text-shadow: 2px 2px 4px rgba(0,0,0,0.25); }
    .main-header .subtitle { color: rgba(255,255,255,0.95); font-size: 1.1rem; margin-top: 0.5rem; font-weight: 500; text-shadow: 1px 1px 2px rgba(0,0,0,0.2); }
    .progress-container { background: white; border-radius: 12px; padding: 1.5rem; margin-bottom: 2rem; box-shadow: 0 4px 20px rgba(0,0,0,0.2); border: 1px solid #e8e8e8; }
    .progress-bar { display: flex; justify-content: space-between; align-items: center; margin-bottom: 1rem; }
    .progress-step { display: flex; flex-direction: column; align-items: center; flex: 1; position: relative; }
    .progress-step:not(:last-child)::after { content: ''; position: absolute; top: 20px; right: -50%; width: 100%; height: 3px; background: #e0e0e0; z-index: 0; }
    .progress-step.active:not(:last-child)::after { background: linear-gradient(90deg, #060c71, #f9d20e); }
    .step-circle { width: 40px; height: 40px; border-radius: 50%; background: #e0e0e0; color: #666; display: flex; align-items: center; justify-content: center; font-weight: 700; font-size: 16px; position: relative; z-index: 1; transition: all 0.3s ease; }
    .step-circle.active { background: linear-gradient(135deg, #060c71, #2a3bb8); color: white; transform: scale(1.1); box-shadow: 0 4px 15px rgba(6, 12, 113, 0.4); }
    .step-circle.completed { background: linear-gradient(135deg, #f9d20e, #ffe34a); color: #060c71; font-weight: 900; }
    .step-label { margin-top: 0.5rem; font-size: 0.9rem; font-weight: 600; color: #666; text-align: center; }
    .step-label.active { color: #060c71; }
    .section-card { background: white; border-radius: 15px; padding: 2rem; margin-bottom: 1.5rem; box-shadow: 0 4px 20px rgba(0,0,0,0.08); border: 1px solid #e8e8e8; transition: all 0.3s ease; }
    .section-card:hover { box-shadow: 0 8px 30px rgba(0,0,0,0.12); transform: translateY(-2px); }
    .section-title { color: #060c71; font-size: 1.4rem; font-weight: 700; margin-bottom: 1rem; padding-bottom: 0.5rem; border-bottom: 2px solid #f9d20e; display: inline-block; }
    .stButton > button {
        background: linear-gradient(135deg, #060c71 0%, #2a3bb8 100%); color: white; border: none; border-radius: 10px;
        padding: 0.75rem 2rem; font-size: 16px; font-weight: 600; transition: all 0.3s ease; box-shadow: 0 4px 15px rgba(6,12,113,0.3); width: 100%;
    }
    .stButton > button:hover { background: linear-gradient(135deg, #f9d20e 0%, #ffe34a 100%); color: #060c71; transform: translateY(-2px); box-shadow: 0 6px 20px rgba(249,210,14,0.4); }
    .stButton > button:disabled { background: #cccccc; color: #666666; transform: none; box-shadow: none; }
    .stTextInput > div > div > input, .stTextArea > div > div > textarea, .stDateInput > div > div > input {
        border-radius: 8px; border: 2px solid #e0e0e0; padding: 0.75rem; font-size: 16px; transition: all 0.3s ease;
    }
    .stTextInput > div > div > input:focus, .stTextArea > div > div > textarea:focus, .stDateInput > div > div > input:focus {
        border-color: #060c71; box-shadow: 0 0 0 3px rgba(6, 12, 113, 0.1);
    }
    .stFileUploader > div { border: 2px dashed #060c71; border-radius: 10px; padding: 1rem; text-align: center; background: rgba(6,12,113,0.02); transition: all 0.3s ease; }
    .stFileUploader > div:hover { background: rgba(6, 12, 113, 0.05); border-color: #f9d20e; }
    .streamlit-expanderHeader { background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%); border-radius: 10px; border: 2px solid #e0e0e0; font-weight: 600; color: #060c71; }
    .streamlit-expanderContent { border: 2px solid #e0e0e0; border-top: none; border-radius: 0 0 10px 10px; background: white; }
    .stDownloadButton > button { background: linear-gradient(135deg, #28a745 0%, #34ce57 100%); color: white; border: none; border-radius: 8px; padding: 0.5rem 1rem; font-weight: 600; transition: all 0.3s ease; width: 100%; }
    .stDownloadButton > button:hover { background: linear-gradient(135deg, #218838 0%, #28a745 100%); transform: translateY(-1px); }
    .stRadio > div { background: white; padding: 1rem; border-radius: 10px; border: 2px solid #e0e0e0; }
    #MainMenu {visibility: hidden;} footer {visibility: hidden;} .stDeployButton {display: none;}
</style>
""", unsafe_allow_html=True)

# Initialize session state
if "step" not in st.session_state:
    st.session_state.step = 1

# Progress Steps Configuration — UPDATED to include Flowchart step
STEPS = [
    {"number": 1, "label": "Project Details", "icon": "📋"},
    {"number": 2, "label": "Cover Letter", "icon": "📝"},
    {"number": 3, "label": "Executive Summary", "icon": "📊"},
    {"number": 4, "label": "Flowchart", "icon": "🗺️"},
    {"number": 5, "label": "Final Proposal", "icon": "📄"}
]

def render_progress_bar(current_step):
    progress_html = '<div class="progress-container"><div class="progress-bar">'
    for step in STEPS:
        is_active = step["number"] == current_step
        is_completed = step["number"] < current_step
        
        circle_class = "step-circle"
        label_class = "step-label"
        
        if is_completed:
            circle_class += " completed"
            content = "✓"
        elif is_active:
            circle_class += " active"
            label_class += " active"
            content = str(step["number"])
        else:
            content = str(step["number"])
        
        step_class = "progress-step"
        if is_active:
            step_class += " active"
        
        progress_html += f'<div class="{step_class}"><div class="{circle_class}">{content}</div><div class="{label_class}">{step["label"]}</div></div>'
    progress_html += '</div></div>'
    st.markdown(progress_html, unsafe_allow_html=True)

def render_header():
    header_col1, header_col2 = st.columns([0.15, 0.85])
    with header_col1:
        try:
            st.image("Input\\Static_AboutCompany\\Falcon-Autotech_Logo.png", width=300)
        except:
            st.write("")
    with header_col2:
        st.markdown('''
        <div class="main-header">
            <h1>Sales Proposal Automation System (SPAS)</h1>
            <div class="subtitle">Boosting proposal efficiency with powerful and intuitive tool</div>
        </div>
        ''', unsafe_allow_html=True)

def render_section_card(title, content_func):
    st.markdown(f'''
    <div class="section-card">
        <div class="section-title">{title}</div>
    </div>
    ''', unsafe_allow_html=True)
    with st.container():
        content_func()
# === NEW: Draw.io / Selenium helpers ===
from dataclasses import dataclass

@dataclass
class _FlowAppState:
    driver: object | None = None
    mode: str | None = None          # "drawio" when editor is open
    download_dir: str = tempfile.mkdtemp(prefix="flowchart_")
    viewer_url: str | None = None    # diagrams.net viewer link (live preview)
    last_refresh_ts: float | None = None

def _ensure_flow_state():
    if "flow_state" not in st.session_state:
        st.session_state.flow_state = _FlowAppState()


class _SeleniumHelper:
    @staticmethod
    def create_driver(headless: bool, download_dir: Optional[str]):
        from selenium import webdriver
        from selenium.webdriver.chrome.service import Service
        from webdriver_manager.chrome import ChromeDriverManager

        options = webdriver.ChromeOptions()
        if headless: options.add_argument("--headless=new")
        options.add_argument("--disable-gpu")
        options.add_argument("--no-sandbox")
        options.add_argument("--disable-dev-shm-usage")
        options.add_argument("--start-maximized")
        options.add_experimental_option("detach", True)
        options.add_experimental_option("excludeSwitches", ["enable-logging"])

        if download_dir:
            os.makedirs(download_dir, exist_ok=True)
            prefs = {
                "download.default_directory": os.path.abspath(download_dir),
                "download.prompt_for_download": False,
                "download.directory_upgrade": True,
                "safebrowsing.enabled": True,
                "profile.default_content_setting_values.automatic_downloads": 1,
            }
            options.add_experimental_option("prefs", prefs)

        driver = webdriver.Chrome(
            service=Service(ChromeDriverManager().install()),
            options=options
        )
        driver.set_page_load_timeout(900)
        driver.implicitly_wait(2)
        return driver

    @staticmethod
    def js_click(driver, element):
        driver.execute_script("arguments[0].scrollIntoView({block:'center'});", element)
        time.sleep(0.05)
        try:
            element.click()
        except Exception:
            driver.execute_script("arguments[0].click();", element)

    @staticmethod
    def wait_and_click(driver, xpath: str, timeout: int = 30):
        from selenium.webdriver.common.by import By
        from selenium.webdriver.support.ui import WebDriverWait
        from selenium.webdriver.support import expected_conditions as EC
        element = WebDriverWait(driver, timeout).until(
            EC.element_to_be_clickable((By.XPATH, xpath))
        )
        _SeleniumHelper.js_click(driver, element)
        return element

    @staticmethod
    def wait_and_type(driver, xpath: str, text: str, clear: bool = True, timeout: int = 30):
        from selenium.webdriver.common.by import By
        from selenium.webdriver.support.ui import WebDriverWait
        from selenium.webdriver.support import expected_conditions as EC
        element = WebDriverWait(driver, timeout).until(
            EC.presence_of_element_located((By.XPATH, xpath))
        )
        driver.execute_script("arguments[0].scrollIntoView({block:'center'});", element)
        if clear:
            try: element.clear()
            except Exception: pass
        try: element.send_keys(text)
        except Exception: pass
        return element

    @staticmethod
    def get_latest_file(directory: str) -> Optional[str]:
        files = sorted(glob.glob(os.path.join(directory, "*")), key=os.path.getmtime)
        return files[-1] if files else None

class _DrawIO:
    @staticmethod
    def open_and_insert(mermaid_code: str, headless: bool, download_dir: str):
        from selenium.webdriver.common.action_chains import ActionChains
        driver = _SeleniumHelper.create_driver(headless=headless, download_dir=download_dir)
        driver.get("https://app.diagrams.net/")
        try:
            ActionChains(driver).pause(0.6).send_keys("\ue00c").perform()  # ESC to close cookie/modals
        except Exception:
            pass

        # Menu: Insert → Advanced → Mermaid
        _SeleniumHelper.wait_and_click(driver, "/html/body/div[1]/div[1]/a[4]", timeout=30)   # Insert
        time.sleep(0.15)
        _SeleniumHelper.wait_and_click(driver, "/html/body/div[9]/table/tbody/tr[13]/td[2]", timeout=30)  # Advanced
        time.sleep(0.15)
        _SeleniumHelper.wait_and_click(driver, "/html/body/div[10]/table/tbody/tr[15]/td[2]", timeout=30)  # Mermaid
        time.sleep(0.15)

        # Paste Mermaid
        _SeleniumHelper.wait_and_type(driver, "/html/body/div[10]/div/textarea", mermaid_code, clear=True, timeout=30)
        time.sleep(0.1)

        # Insert
        _SeleniumHelper.wait_and_click(driver, "/html/body/div[10]/div/div/button[2]", timeout=30)
        time.sleep(1.0)
        return driver
        
    @staticmethod
    def export_xml(driver, download_dir: str, timeout_sec: int = 60) -> str:
        """
        Robust XML export for diagrams.net.
        Attempt A: Direct JS export via ui.editor.getGraphXml()  (no UI clicks).
        Attempt B: Hover-driven menu flow using text-based locators (no brittle XPaths).
        Returns the file path of the saved .xml/.drawio
        """
        import os, time
        from pathlib import Path
        from selenium.webdriver.common.by import By
        from selenium.webdriver.common.action_chains import ActionChains
        from selenium.webdriver.support.ui import WebDriverWait
        from selenium.webdriver.support import expected_conditions as EC
        from selenium.common.exceptions import TimeoutException

        wait = WebDriverWait(driver, 20)
        Path(download_dir).mkdir(parents=True, exist_ok=True)

        def _switch_to_editor_iframe():
            driver.switch_to.default_content()
            # Try common selectors; fall back to first visible iframe
            candidates = driver.find_elements(By.CSS_SELECTOR,
                "iframe[src*='diagrams.net'],iframe[src*='embed.diagrams.net'],iframe[id*='geEditor'],iframe")
            for fr in candidates:
                if fr.is_displayed():
                    try:
                        driver.switch_to.frame(fr)
                        # ensure UI ready
                        ok = WebDriverWait(driver, 5).until(
                            lambda d: d.execute_script("return !!(window.ui && ui.editor && ui.editor.getGraphXml)")
                        )
                        if ok:
                            return
                    except Exception:
                        driver.switch_to.default_content()
            # If not in iframe, try top-level
            driver.switch_to.default_content()
            try:
                WebDriverWait(driver, 5).until(
                    lambda d: d.execute_script("return !!(window.ui && ui.editor && ui.editor.getGraphXml)")
                )
            except Exception:
                pass

        def _write_xml(xml_text: str) -> str:
            ts = time.strftime("%Y%m%d_%H%M%S")
            fname = f"diagram_{ts}.drawio"
            fpath = os.path.join(download_dir, fname)
            with open(fpath, "w", encoding="utf-8") as f:
                f.write(xml_text)
            return fpath

        def _last_visible_popup_xpath(text_ci: str) -> str:
            return ("(//div[contains(@class,'mxPopupMenu') and not(contains(@style,'display: none'))])[last()]"
                    "//*[self::td or self::div or self::span]"
                    f"[contains(translate(normalize-space(.),'ABCDEFGHIJKLMNOPQRSTUVWXYZ','abcdefghijklmnopqrstuvwxyz'),'{text_ci.lower()}')]")

        def _last_visible_dialog_button_xpath(*labels) -> str:
            ors = " or ".join([
                f"contains(translate(normalize-space(.),'ABCDEFGHIJKLMNOPQRSTUVWXYZ','abcdefghijklmnopqrstuvwxyz'),'{l.lower()}')"
                for l in labels
            ])
            return ("(//div[contains(@class,'geDialog') and not(contains(@style,'display: none'))])[last()]"
                    f"//button[{ors}]")

        def _hover_click(el):
            ActionChains(driver).move_to_element(el).pause(0.1).click().perform()

        # ---------- Attempt A: direct JS export (most reliable, no UI clicks) ----------
        _switch_to_editor_iframe()
        try:
            xml = driver.execute_script(
                "return (window.ui && window.mxUtils && ui.editor) ? "
                "mxUtils.getXml(ui.editor.getGraphXml()) : null;"
            )
            if isinstance(xml, str) and len(xml) > 100 and ("<mxfile" in xml or "<mxGraphModel" in xml):
                return _write_xml(xml)
        except Exception:
            pass  # fall through to Attempt B

        # ---------- Attempt B: hover-driven UI flow, text selectors ----------
        # Close any overlaying dialogs/masks that can block clicks
        try:
            driver.execute_script("""
                (function(){
                var masks = document.querySelectorAll('.geMask');
                masks.forEach(m => { if (m.offsetParent) m.click?.(); });
                })();
            """)
        except Exception:
            pass

        # 1) File (menubar)
        file_btn = wait.until(EC.element_to_be_clickable((
            By.XPATH,
            "//div[contains(@class,'geMenubar')]//*[self::a or self::div or self::span]"
            "[translate(normalize-space(.),'ABCDEFGHIJKLMNOPQRSTUVWXYZ','abcdefghijklmnopqrstuvwxyz')='file']"
        )))
        _hover_click(file_btn)
        time.sleep(0.15)

        # 2) Export As (or 'Export as')
        try:
            export_as = wait.until(EC.element_to_be_clickable((By.XPATH, _last_visible_popup_xpath("export as"))))
        except TimeoutException:
            export_as = wait.until(EC.element_to_be_clickable((By.XPATH, _last_visible_popup_xpath("export as"))))
        _hover_click(export_as)
        time.sleep(0.15)

        # 3) XML / XML… (ellipsis variants)
        try:
            xml_item = wait.until(EC.element_to_be_clickable((By.XPATH, _last_visible_popup_xpath("xml"))))
        except TimeoutException:
            # try with ellipsis in text
            xml_item = wait.until(EC.element_to_be_clickable((By.XPATH, _last_visible_popup_xpath("xml…"))))
        _hover_click(xml_item)

        # 4) Export (dialog)
        export_btn = wait.until(EC.element_to_be_clickable((By.XPATH, _last_visible_dialog_button_xpath("Export"))))
        _hover_click(export_btn)

        # 5) Save/Download (dialog)
        try:
            save_btn = wait.until(EC.element_to_be_clickable((By.XPATH, _last_visible_dialog_button_xpath("Download", "Save"))))
        except TimeoutException:
            save_btn = wait.until(EC.element_to_be_clickable((By.XPATH, _last_visible_dialog_button_xpath("Export"))))
        _hover_click(save_btn)

        # Wait for a file to appear in download_dir (non-.crdownload, minimal size)
        before = _SeleniumHelper.get_latest_file(download_dir)
        end = time.time() + timeout_sec
        while time.time() < end:
            candidate = _SeleniumHelper.get_latest_file(download_dir)
            if candidate and (before is None or os.path.normpath(candidate) != os.path.normpath(before)):
                if not candidate.endswith(".crdownload"):
                    try:
                        if os.path.getsize(candidate) > 2048:
                            return candidate
                    except Exception:
                        return candidate
            time.sleep(0.3)

        raise RuntimeError("XML export timeout (no file downloaded)")





    @staticmethod
    def export_png(driver, download_dir: str, timeout_sec: int = 900) -> Optional[str]:
        from selenium.webdriver.common.by import By
        from selenium.webdriver.support.ui import WebDriverWait
        from selenium.webdriver.support import expected_conditions as EC

        wait = WebDriverWait(driver, 10)
        before = _SeleniumHelper.get_latest_file(download_dir)
        try:
            file_btn = wait.until(EC.element_to_be_clickable((By.XPATH, "/html/body/div[1]/div[1]/a[1]")))
            file_btn.click(); time.sleep(0.25)
            export_as = _SeleniumHelper.wait_and_click(By.XPATH, "/html/body/div[9]/table/tbody/tr[14]/td[2]")  # Export As
            time.sleep(0.25)
            png_item = wait.until(EC.element_to_be_clickable((driver, "//div[contains(@class,'mxPopupMenu')][last()]//td[normalize-space()='PNG']")))
            _SeleniumHelper.js_click(driver, png_item)
            _SeleniumHelper.wait_and_click(By.XPATH, "/html/body/div[10]/div[1]/div[2]/button[2]", timeout=20)  # Export
            _SeleniumHelper.wait_and_click(By.XPATH, "/html/body/div[10]/div/div[2]/button[3]", timeout=20)      # Download

            end_time = time.time() + timeout_sec
            while time.time() < end_time:
                candidate = _SeleniumHelper.get_latest_file(download_dir)
                if candidate and (before is None or os.path.normpath(candidate) != os.path.normpath(before)):
                    if candidate.lower().endswith(".png") and not candidate.endswith(".crdownload"):
                        if os.path.getsize(candidate) > 1024:
                            return candidate
                time.sleep(0.3)
        except Exception:
            return None
        return None

    @staticmethod
    def xml_to_viewer_url(xml_str: str) -> str:
        raw = xml_str.encode("utf-8")
        comp = zlib.compressobj(level=9, wbits=-15)
        deflated = comp.compress(raw) + comp.flush()
        b64 = base64.b64encode(deflated).decode("ascii")
        return f"https://viewer.diagrams.net/?lightbox=1&nav=1&layers=1&noSaveBtn=1#R{b64}"

# --- ALWAYS rebuild the flowchart one-pager from the latest diagram before merge ---
def _valid_png(png: bytes) -> Optional[bytes]:
    try:
        from PIL import Image as PILImage
        im = PILImage.open(io.BytesIO(png))
        im.load()  # force decode
        if im.width < 8 or im.height < 8:
            return None
        # Flatten/crop to avoid alpha/huge white borders & Office quirks
        return _auto_crop_png_whitespace(png)
    except Exception:
        return None

    try:
        latest_png = None
        # Prefer a fresh export if diagrams.net editor is active
        try:
            latest_png = _export_latest_flowchart_png()
        except Exception:
            latest_png = None

        if not latest_png:
            # fallback: whatever we captured at Step 4, then any previous render
            latest_png = (st.session_state.get("flowchart_final_png_bytes")
                        or st.session_state.get("mermaid_png_bytes"))

        latest_png = _valid_png(latest_png) if latest_png else None


        st.session_state.flowchart_docx_buffer = make_docx_fit_one_page_from_png_or_code(
            latest_png,
            st.session_state.get("mermaid_code", ""),
            title="Concept Description",
            client_logo=st.session_state.get("client_logo"),
            falcon_logo=st.session_state.get("falcon_logo"),
            client_name=st.session_state.get("client_name", ""),
            project_name=st.session_state.get("project_title", ""),
        )
    except Exception as e:
        st.error(f"Could not rebuild flowchart page: {e}")
        st.stop()
def ensure_falcon_section_title_style(doc,
                                      style_name="FalconSectionTitle",
                                      font_name="Calibri",
                                      font_size_pt=16,
                                      rgb=(2, 12, 115),  # #004AAD
                                      underline=True,
                                      bold=True):
    """Create (or reuse) a consistent section title style across all pages."""
    styles = doc.styles
    try:
        style = styles[style_name]
    except KeyError:
        style = styles.add_style(style_name, WD_STYLE_TYPE.PARAGRAPH)
        # paragraph-level defaults
        style.paragraph_format.space_before = Pt(0)
        style.paragraph_format.space_after = Pt(6)
        style.paragraph_format.left_indent = Pt(0)
        # font
        f = style.font
        f.name = font_name
        f.size = Pt(font_size_pt)
        f.bold = bold
        f.color.rgb = RGBColor(*rgb)
    return style_name

def _require_valid_png(png_bytes: bytes) -> bytes:
    """Decode + sanity check + crop + flatten. Raise if not usable."""
    if not png_bytes or len(png_bytes) < 100:  # Changed from checking dimensions
        raise ValueError("Flowchart PNG missing or too small.")
    
    try:
        im = PILImage.open(io.BytesIO(png_bytes))
        im.load()
    except Exception as e:
        raise ValueError(f"Flowchart PNG is corrupt: {e}")

    # Remove the 8x8 dimension check - allow any valid decoded image
    
    # crop transparent/near-white borders
    bbox = None
    if im.mode == "RGBA":
        bbox = im.getchannel("A").getbbox()
    if not bbox:
        gray = ImageOps.grayscale(im.convert("RGB"))
        bg = gray.point(lambda p: 255 if p > 245 else 0)
        bbox = ImageOps.invert(bg).getbbox()
    if bbox:
        im = im.crop(bbox)

    # Ensure minimum size AFTER cropping
    if im.width < 50 or im.height < 50:
        raise ValueError(f"Flowchart image too small after cropping ({im.width}x{im.height})")

    # flatten alpha onto white
    if im.mode == "RGBA":
        bg = PILImage.new("RGB", im.size, (255, 255, 255))
        bg.paste(im, mask=im.split()[-1])
        im = bg

    out = io.BytesIO()
    im.save(out, format="PNG", optimize=True)
    out.seek(0)
    return out.getvalue()

def _docx_from_png_onepage(png_bytes: bytes,
                           title: str,
                           client_logo=None, falcon_logo=None,
                           client_name="", project_name="") -> io.BytesIO:
    """Build an A4 portrait one-pager with the PNG filling the printable area."""
    doc = Document()
    # headers/footers if you already have a helper:
    try:
        add_page_headers(doc, client_logo, falcon_logo, client_name, project_name)
    except Exception:
        pass

    # Force A4 portrait (inches)
    section = doc.sections[0]
    section.page_width  = Inches(8.27)
    section.page_height = Inches(11.69)
    # Leave whatever margins your header/footer logic expects

    # Title
    if title:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
        title_style = ensure_falcon_section_title_style(doc)
        p = doc.add_paragraph(title, style=title_style)


    # Compute max width within current margins
    left = float(section.left_margin.inches)
    right = float(section.right_margin.inches)
    top = float(section.top_margin.inches)
    bottom = float(section.bottom_margin.inches)

    max_w = 8.27 - left - right
    # Leave a little vertical space for title; scale by width and let Word keep aspect
    image_stream = io.BytesIO(png_bytes)
    doc.add_paragraph("")  # spacer
    run = doc.add_paragraph().add_run()
    run.add_picture(image_stream, width=Inches(max_w))

    buf = io.BytesIO()
    doc.save(buf); buf.seek(0)
    return buf

def _auto_refresh_drawio_preview() -> bool:
    """
    Export XML from the open Draw.io tab and update the right-column preview iframe.
    Returns True if preview was updated successfully.
    """
    _ensure_flow_state()
    fs = st.session_state.flow_state
    
    if fs.mode != "drawio" or fs.driver is None:
        st.warning("Draw.io editor is not active. Please click 'Edit' first.")
        return False
    
    try:
        # 1. Check if browser session is still alive
        try:
            _ = fs.driver.current_url  # Ping the driver
        except Exception as e:
            st.error("⚠️ Draw.io browser session lost. Please click 'Edit' to reopen.")
            fs.driver = None
            fs.mode = None
            return False
        
        # 2. Check if user is still on diagrams.net domain
        current_url = fs.driver.current_url
        if "diagrams.net" not in current_url and "draw.io" not in current_url:
            st.warning("⚠️ Browser navigated away from Draw.io. Please return to the diagram.")
            return False
        
        # 3. Export XML with extended timeout and retry logic
        max_retries = 2
        xml_file = None
        last_error = None
        
        for attempt in range(max_retries):
            try:
                # Increase timeout for complex diagrams
                xml_file = _DrawIO.export_xml(
                    fs.driver, 
                    fs.download_dir, 
                    timeout_sec=900  # Increased from 60
                )
                break  # Success
            except Exception as e:
                last_error = e
                if attempt < max_retries - 1:
                    st.info(f"Export attempt {attempt + 1} failed, retrying...")
                    time.sleep(1.5)
                continue
        
        if not xml_file:
            raise RuntimeError(f"XML export failed after {max_retries} attempts: {last_error}")
        
        # 4. Read and validate XML
        if not os.path.exists(xml_file):
            raise FileNotFoundError(f"Exported XML file not found: {xml_file}")
        
        xml_text = pathlib.Path(xml_file).read_text(encoding="utf-8", errors="ignore")
        
        if not xml_text or len(xml_text) < 50:
            raise ValueError("Exported XML is empty or too small")
        
        # 5. Generate viewer URL
        fs.viewer_url = _DrawIO.xml_to_viewer_url(xml_text)
        fs.last_refresh_ts = time.time()
        
        return True
        
    except FileNotFoundError as e:
        st.error(f"❌ Export file not found: {e}")
        return False
    
    except ValueError as e:
        st.error(f"❌ Invalid diagram data: {e}")
        return False
    
    except RuntimeError as e:
        st.error(f"❌ Export timeout: {e}")
        st.info("💡 Tip: Try simplifying your diagram or waiting a few seconds before refreshing.")
        return False
    
    except Exception as e:
        error_msg = str(e)
        
        # Categorize common errors with helpful messages
        if "no such window" in error_msg.lower() or "invalid session" in error_msg.lower():
            st.error("❌ Draw.io browser window was closed. Please click 'Edit' to reopen.")
            fs.driver = None
            fs.mode = None
        elif "timeout" in error_msg.lower():
            st.error("❌ Export timed out. The diagram may be too complex or Draw.io is busy.")
            st.info("💡 Try: Wait 5 seconds and click refresh again, or simplify your diagram.")
        elif "element not found" in error_msg.lower():
            st.error("❌ Could not locate export menu in Draw.io. The page may have changed.")
        else:
            st.error(f"❌ Refresh failed: {error_msg}")
        
        return False

def _auto_crop_png_whitespace(png_bytes: bytes) -> bytes:
    """
    Crop transparent/near-white borders and flatten to RGB to avoid Word/Office
    placeholder issues with large-alpha PNGs.
    """
    try:
        im = PILImage.open(io.BytesIO(png_bytes))
        if im.mode not in ("RGB", "RGBA"):
            im = im.convert("RGBA")

        # 1) find bounding box
        bbox = None
        if im.mode == "RGBA":
            alpha = im.getchannel("A")
            bbox = alpha.getbbox()

        if not bbox:
            # near-white crop for opaque images/screenshots
            gray = ImageOps.grayscale(im.convert("RGB"))
            # consider anything > 245 as "white"
            bg = gray.point(lambda p: 255 if p > 245 else 0)
            bbox = ImageOps.invert(bg).getbbox()

        if bbox:
            im = im.crop(bbox)

        # 2) flatten alpha onto white to avoid Office rendering quirks
        if im.mode == "RGBA":
            bg = PILImage.new("RGB", im.size, (255, 255, 255))
            bg.paste(im, mask=im.split()[-1])
            im = bg

        out = io.BytesIO()
        # keep PNG (now RGB, no alpha) which python-docx embeds cleanly
        im.save(out, format="PNG", optimize=True)
        out.seek(0)
        return out.getvalue()
    except Exception:
        return png_bytes

def _indent_exec_summary_bullets(es_buf: io.BytesIO,
                                 one_indent_in=0.25,
                                 base_left_in=0.5,
                                 hang_in=0.25) -> io.BytesIO:
    """
    Increase bullet indent by one 'tab stop' (≈0.25") reliably.
    1) Prefer Word's built-in 'List Bullet 2' (one level deeper) if available.
    2) Otherwise, keep the bullet glyph and apply manual left + hanging indent.
    """
    es_buf.seek(0)
    doc = Document(es_buf)

    def _try_style_list_bullet_2(p):
        """Try to switch to List Bullet 2 (one level deeper)."""
        try:
            p.style = doc.styles["List Bullet 2"]
            return True
        except Exception:
            return False

    for p in doc.paragraphs:
        text = (p.text or "")
        m = BULLET_RX.match(text)
        if not m:
            continue

        glyph, content = m.groups()
        # Normalize paragraph text to content only; styling will add the bullet
        p.clear()
        r = p.add_run(content.strip())
        r.font.size = r.font.size or Pt(11)

        # Try deepening the bullet using Word’s built-in style first
        if _try_style_list_bullet_2(p):
            # (Optional) tiny nudge further right if you want it obvious)
            pf = p.paragraph_format
            pf.left_indent = (pf.left_indent or Inches(0)) + Inches(one_indent_in)
            # Keep hanging so wrapped lines align nicely
            pf.first_line_indent = Inches(-hang_in)
        else:
            # Fallback: manual bullet + indent
            # Rebuild as "•\tcontent" and push right by base_left_in + one_indent_in
            p.clear()
            rb = p.add_run("•\t")
            rb.font.size = Pt(11)
            rc = p.add_run(content.strip())
            rc.font.size = Pt(11)

            pf = p.paragraph_format
            pf.left_indent = Inches(base_left_in + one_indent_in)   # shift right
            pf.first_line_indent = Inches(-hang_in)                 # hanging
            # We avoid style-numbering entirely; this renders consistently in Word

        # Extra safety: if paragraph already had numbering, bump ilvl by +1
        try:
            pPr = p._p.get_or_add_pPr()
            numPr = pPr.find(qn('w:numPr'))
            if numPr is not None:
                ilvl = numPr.find(qn('w:ilvl'))
                if ilvl is None:
                    ilvl = OxmlElement('w:ilvl')
                    numPr.append(ilvl)
                # increment level by 1 (cap at 8)
                current = int(ilvl.get(qn('w:val'), '0'))
                ilvl.set(qn('w:val'), str(min(current + 1, 8)))
        except Exception:
            pass

    out = io.BytesIO()
    doc.save(out)
    out.seek(0)
    return out

def _export_latest_flowchart_png() -> Optional[bytes]:
    """
    If Draw.io is active, export latest PNG from it (fallback to viewer screenshot).
    Otherwise render PNG via Kroki from current mermaid_code.
    Returns PNG bytes or None.
    """
    _ensure_flow_state()
    fs = st.session_state.flow_state

    # Draw.io route
    if fs.mode == "drawio" and fs.driver is not None:
        # try native PNG export
        png_path = _DrawIO.export_png(fs.driver, fs.download_dir, timeout_sec=900)
        if png_path and os.path.exists(png_path):
            fs.last_png_path = png_path
            return open(png_path, "rb").read()

        # fallback: export XML + open viewer and screenshot it headless
        try:
            xml_file = _DrawIO.export_xml(fs.driver, fs.download_dir, timeout_sec=900)
            xml_text = pathlib.Path(xml_file).read_text(encoding="utf-8", errors="ignore")
            viewer_url = _DrawIO.xml_to_viewer_url(xml_text)

            drv = _SeleniumHelper.create_driver(headless=True, download_dir=None)
            try:
                drv.set_window_size(1920, 1200)
                drv.get(viewer_url)
                time.sleep(2.0)
                png_bytes = drv.get_screenshot_as_png()
                fs.last_png_path = os.path.join(fs.download_dir, f"drawio_view_{int(time.time())}.png")
                Image.open(io.BytesIO(png_bytes)).save(fs.last_png_path)
                return open(fs.last_png_path, "rb").read()
            finally:
                try: drv.quit()
                except Exception: pass
        except Exception:
            pass

    # Plain Mermaid route
    code = st.session_state.get("mermaid_code", "")
    if code.strip():
        png = mermaid_to_png_via_kroki(code)
        return png
    return None


# Main Header
render_header()
# Progress Bar
render_progress_bar(st.session_state.step)

# ========================
# STEP 1: Project Details
# ========================
if st.session_state.step == 1:
    def step1_content():
        st.markdown("#### Client Information")
        col1, col2 = st.columns(2)
        
        with col1:
            client_name = st.text_input(
                "Client Name*", 
                value=st.session_state.get("client_name", ""),
                help="Enter the full name of the client organization"
            )
            project_title = st.text_input(
                "Project Title*", 
                value=st.session_state.get("project_title", ""),
                help="Enter the complete project title"
            )
            offer_ref = st.text_input(
                "Offer Reference Number*", 
                value=st.session_state.get("offer_ref", ""),
                help="Enter the official reference number for this proposal"
            )
        
        with col2:
            executives = st.text_area(
                "Key Executives (comma separated)", 
                value=st.session_state.get("executives", ""),
                height=100,
                help="List key client executives"
            )
            invitation_date = st.date_input(
                "Invitation Date (Optional)", 
                value=st.session_state.get("invitation_date", None),
                help="Date when the invitation was received"
            )
            meeting_date = st.date_input(
                "Meeting Date (Optional)", 
                value=st.session_state.get("meeting_date", None),
                help="Scheduled meeting or presentation date"
            )
        
        st.markdown("---")
        st.markdown("#### Document Assets")
        
        upload_col1, upload_col2, upload_col3 = st.columns(3)
        
        with upload_col1:
            st.markdown("**Solution Design Document**")
            uploaded_pdf = st.file_uploader(
                "Upload PDF", 
                type=["pdf"], 
                key="pdf_upl",
                help="Upload the technical solution design document"
            )
        
        with upload_col2:
            st.markdown("**Client Logo**")
            uploaded_client_logo = st.file_uploader(
                "Upload Logo", 
                type=["png", "jpg", "jpeg"], 
                key="logo_upl",
                help="Upload the client's official logo"
            )
        
        with upload_col3:
            st.markdown("**RFQ Diagram**")
            rfq_img = st.file_uploader(
                "Upload Diagram", 
                type=["png", "jpg", "jpeg"], 
                key="rfq_upl",
                help="Upload the RFQ architecture or process diagram"
            )
        
        if uploaded_client_logo:
            try:
                st.session_state.client_logo = uploaded_client_logo.getvalue()
                st.success("Client logo uploaded successfully")
            except Exception as e:
                st.error(f"Unable to read client logo: {e}")

        if uploaded_pdf:
            try:
                pdf_bytes = uploaded_pdf.getvalue()
                with open("temp_uploaded.pdf", "wb") as f:
                    f.write(pdf_bytes)
                st.session_state["solution_pdf_bytes"] = pdf_bytes
                st.success("Solution design PDF uploaded successfully")
            except Exception as e:
                st.error(f"Unable to cache PDF: {e}")

        if rfq_img and client_name and project_title and offer_ref:
            try:
                img_data = rfq_img.read()
                rfq_doc = Document()
                add_page_headers(
                    rfq_doc,
                    st.session_state.get("client_logo"),
                    st.session_state.get("falcon_logo"),
                    client_name,
                    project_title,
                )
                add_response_heading(rfq_doc, f"Response to Requirement for {project_title}")
                rfq_doc.add_paragraph("")
                rfq_doc.add_picture(io.BytesIO(img_data), width=Inches(6))
                rfq_doc.add_paragraph("")
                add_response_heading(rfq_doc, f"Proposal Reference – {offer_ref}")
                rfq_doc.add_paragraph("")
                add_contact_box(rfq_doc)

                rfq_buf = io.BytesIO()
                rfq_doc.save(rfq_buf)
                rfq_buf.seek(0)
                st.session_state.manual_rfq_buffer = rfq_buf
                st.success("RFQ diagram integrated into proposal structure")
            except Exception as e:
                st.error(f"Failed to build RFQ document: {e}")
        
        st.markdown("---")
        
        required_fields = [client_name, project_title, offer_ref]
        missing_fields = []
        if not client_name: missing_fields.append("Client Name")
        if not project_title: missing_fields.append("Project Title")
        if not offer_ref: missing_fields.append("Offer Reference")
        
        if missing_fields:
            st.warning(f"Please fill in the following required fields: {', '.join(missing_fields)}")
        
        col1, col2, col3 = st.columns([1, 2, 1])
        with col1:
            proceed_button = st.button(
                "Next →",
                use_container_width=True,
                disabled=bool(missing_fields),
                key="proceed_step1"
            )
        
        if proceed_button:
            st.session_state.client_name = client_name
            st.session_state.executives = executives
            st.session_state.offer_ref = offer_ref
            st.session_state.project_title = project_title
            st.session_state.invitation_date = invitation_date
            st.session_state.meeting_date = meeting_date
            st.session_state.step = 2
            st.rerun()
    
    render_section_card("Project Setup & Asset Upload", step1_content)

# ========================
# STEP 2: Cover Letters
# ========================
elif st.session_state.step == 2:
    def step2_content():
        st.markdown("#### Generate Professional Cover Letters")
        st.markdown("Create personalized cover letters for your proposal using AI-powered content generation.")
        
        col1, col2, col3 = st.columns([1, 2, 1])
        with col1:
            generate_button = st.button(
                "Generate Cover Letter Options",
                key="generate_cover_letters"
            )
        
        if generate_button:
            with st.spinner("Generating personalized cover letters...", show_time=True):
                try:
                    if "generate_cover_letters_concurrent" in globals():
                        cl1, cl2 = generate_cover_letters_concurrent(
                            sender_name="Sandeep Bansal",
                            sender_title="Chief Business Officer",
                            client_name=st.session_state.get("client_name", ""),
                            executives=st.session_state.get("executives", ""),
                            offer_ref=st.session_state.get("offer_ref", ""),
                            project_title=st.session_state.get("project_title", ""),
                            invitation_date=st.session_state.get("invitation_date").strftime("%d-%m-%Y")
                                if st.session_state.get("invitation_date") else None,
                            meeting_date=st.session_state.get("meeting_date").strftime("%d-%m-%Y")
                                if st.session_state.get("meeting_date") else None,
                        )
                    else:
                        cl1 = generate_cover_letter(
                            sender_name="Sandeep Bansal",
                            sender_title="Chief Business Officer",
                            client_name=st.session_state.get("client_name", ""),
                            executives=st.session_state.get("executives", ""),
                            offer_ref=st.session_state.get("offer_ref", ""),
                            project_title=st.session_state.get("project_title", ""),
                            invitation_date=st.session_state.get("invitation_date").strftime("%d-%m-%Y")
                                if st.session_state.get("invitation_date") else None,
                            meeting_date=st.session_state.get("meeting_date").strftime("%d-%m-%Y")
                                if st.session_state.get("meeting_date") else None,
                            letter_date=None,
                        )
                        cl2 = generate_cover_letter(
                            sender_name="Sandeep Bansal",
                            sender_title="Chief Business Officer",
                            client_name=st.session_state.get("client_name", ""),
                            executives=st.session_state.get("executives", ""),
                            offer_ref=st.session_state.get("offer_ref", ""),
                            project_title=st.session_state.get("project_title", ""),
                            invitation_date=st.session_state.get("invitation_date").strftime("%d-%m-%Y")
                                if st.session_state.get("invitation_date") else None,
                            meeting_date=st.session_state.get("meeting_date").strftime("%d-%m-%Y")
                                if st.session_state.get("meeting_date") else None,
                            letter_date=None,
                        )

                    b1 = save_to_docx(
                        cl1,
                        "",
                        st.session_state.get("client_logo"),
                        st.session_state.get("falcon_logo"),
                        st.session_state.get("client_name", ""),
                        st.session_state.get("project_title", "")
                    )
                    b2 = save_to_docx(
                        cl2,
                        "",
                        st.session_state.get("client_logo"),
                        st.session_state.get("falcon_logo"),
                        st.session_state.get("client_name", ""),
                        st.session_state.get("project_title", "")
                    )
                    st.session_state.cover_letters = [{"text": cl1, "buffer": b1}, {"text": cl2, "buffer": b2}]
                    st.success("Cover letters generated successfully!")
                except Exception as e:
                    st.error(f"Failed to generate cover letters: {e}")

        if "cover_letters" in st.session_state:
            st.markdown("#### Review and Select Cover Letter")
            col_cov1, col_cov2 = st.columns([1, 1], gap="medium")
            with col_cov1:
                with st.expander("Cover Letter Option 1",expanded=True):
                    edited1 = st.text_area(
                        "Review and edit if needed:", 
                        st.session_state.cover_letters[0]["text"], 
                        height=300, 
                        key="cl_edit1"
                    )
                    try:
                        st.session_state.cover_letters[0]["buffer"] = save_to_docx(
                            edited1,
                            "",
                            st.session_state.get("client_logo"),
                            st.session_state.get("falcon_logo"),
                            st.session_state.get("client_name", ""),
                            st.session_state.get("project_title", "")
                        )
                    except Exception as e:
                        st.error(f"Failed to rebuild Option 1 DOCX: {e}")
                    
                    st.download_button(
                        "Download Option 1",
                        st.session_state.cover_letters[0]["buffer"],
                        file_name=f'{st.session_state.get("client_name","")}_Cover_Letter_1.docx',
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True
                    )
            with col_cov2:
                with st.expander("Cover Letter Option 2", expanded=True):
                    edited2 = st.text_area(
                        "Review and edit if needed:", 
                        st.session_state.cover_letters[1]["text"], 
                        height=300, 
                        key="cl_edit2"
                    )
                    try:
                        st.session_state.cover_letters[1]["buffer"] = save_to_docx(
                            edited2,
                            "",
                            st.session_state.get("client_logo"),
                            st.session_state.get("falcon_logo"),
                            st.session_state.get("client_name", ""),
                            st.session_state.get("project_title", "")
                        )
                    except Exception as e:
                        st.error(f"Failed to rebuild Option 2 DOCX: {e}")
                    
                    st.download_button(
                        "Download Option 2",
                        st.session_state.cover_letters[1]["buffer"],
                        file_name=f'{st.session_state.get("client_name","")}_Cover_Letter_2.docx',
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True
                    )
            
            col_bt4, col_bt5, col_bt6 = st.columns([1, 2, 1])
            with col_bt4 :
                st.markdown("\n")
                st.markdown("\n")
                st.markdown("#### Select to merge with final Proposal ➠")
            with col_bt5:
                choice = st.radio(
                    " ",
                    ["Option 1", "Option 2"],
                    horizontal=True,
                    key="cover_letter_selection"
                )
                st.session_state.selected_cover_buffer = st.session_state.cover_letters[0 if choice == "Option 1" else 1]["buffer"]

            st.markdown("---")
            
            col1, col2, col3 = st.columns([1, 2, 1])
            with col1:
                continue_button = st.button(
                    "Next →",
                    use_container_width=True,
                    disabled=("selected_cover_buffer" not in st.session_state),
                    key="continue_to_step3"
                )
                if continue_button:
                    st.session_state.step = 3
                    st.rerun()
            with col3:
                if st.button("← Back", key="retry_step2", use_container_width=True):
                    st.session_state.step = 1
                    st.rerun()
    render_section_card("Cover Letter Generation", step2_content)

# ==============================
# STEP 3: Executive Summaries
# ==============================
elif st.session_state.step == 3:
    def step3_content():
        st.markdown("#### Generate Executive Summary")
        st.markdown("Create compelling executive summaries based on your solution design document.")
        
        col1, col2, col3 = st.columns([1, 2, 1])
        with col1:
            generate_summary_button = st.button(
                "Generate Executive Summary Options",
                key="generate_exec_summary"
            )
        
        if generate_summary_button:
            if not os.path.exists("temp_uploaded.pdf"):
                st.error("Please upload the Solution Design PDF in Step 1.")
            else:
                with st.spinner("Analyzing solution document and generating executive summaries...", show_time=True):
                    try:
                        s1 = generate_executive_summary("temp_uploaded.pdf", st.session_state.get("client_name", ""), st.session_state.get("project_title", ""))
                        s2 = generate_executive_summary("temp_uploaded.pdf", st.session_state.get("client_name", ""), st.session_state.get("project_title", ""))

                        b1 = save_to_docx(
                            s1,
                            "Executive Summary",
                            st.session_state.get("client_logo"),
                            st.session_state.get("falcon_logo"),
                            st.session_state.get("client_name", ""),
                            st.session_state.get("project_title", "")
                        )
                        b2 = save_to_docx(
                            s2,
                            "Executive Summary",
                            st.session_state.get("client_logo"),
                            st.session_state.get("falcon_logo"),
                            st.session_state.get("client_name", ""),
                            st.session_state.get("project_title", "")
                        )
                        st.session_state.summaries = [{"text": s1, "buffer": b1}, {"text": s2, "buffer": b2}]
                        st.success("Executive summaries generated successfully!")
                    except Exception as e:
                        st.error(f"Failed to generate executive summaries: {e}")

        if "summaries" in st.session_state:
            st.markdown("#### Review and Select Executive Summary")
            col_exe1, col_exe2 = st.columns([1, 1], gap="medium")
            with col_exe1:
                with st.expander("Executive Summary Option 1",expanded=True):
                    es1 = st.text_area(
                        "Review and edit if needed:", 
                        st.session_state.summaries[0]["text"], 
                        height=300, 
                        key="es_edit1"
                    )
                    try:
                        st.session_state.summaries[0]["buffer"] = save_to_docx(
                            es1,
                            "Executive Summary",
                            st.session_state.get("client_logo"),
                            st.session_state.get("falcon_logo"),
                            st.session_state.get("client_name", ""),
                            st.session_state.get("project_title", "")
                        )
                    except Exception as e:
                        st.error(f"Failed to rebuild Option 1 DOCX: {e}")
                    
                    st.download_button(
                        "Download Option 1",
                        st.session_state.summaries[0]["buffer"],
                        file_name=f'{st.session_state.get("client_name","")}_Executive_Summary_1.docx',
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True
                    )
            with col_exe2:
                with st.expander("Executive Summary Option 2",expanded=True):
                    es2 = st.text_area(
                        "Review and edit if needed:", 
                        st.session_state.summaries[1]["text"], 
                        height=300, 
                        key="es_edit2"
                    )
                    try:
                        st.session_state.summaries[1]["buffer"] = save_to_docx(
                            es2,
                            "Executive Summary",
                            st.session_state.get("client_logo"),
                            st.session_state.get("falcon_logo"),
                            st.session_state.get("client_name", ""),
                            st.session_state.get("project_title", "")
                        )
                    except Exception as e:
                        st.error(f"Failed to rebuild Option 2 DOCX: {e}")
                    
                    st.download_button(
                        "Download Option 2",
                        st.session_state.summaries[1]["buffer"],
                        file_name=f'{st.session_state.get("client_name","")}_Executive_Summary_2.docx',
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True
                    )
            col_bt1, col_bt2, col_bt3 = st.columns([1, 2, 1])
            with col_bt1 :
                st.markdown("\n")
                st.markdown("\n")
                st.markdown("#### Select to merge with final Proposal ➠")
            with col_bt2:
                sel = st.radio(
                    " ",
                    ["Option 1", "Option 2"],
                    horizontal=True,
                    key="es_choice"
                )
                st.session_state.selected_summary_buffer = st.session_state.summaries[0 if sel == "Option 1" else 1]["buffer"]
            
            st.markdown("---")

            c1, c2, c3 = st.columns([1, 2, 1])
            with c1:
                go_flow = st.button("Next →", key="go_flow_btn", use_container_width=True)
                if go_flow:
                    st.session_state.step = 4
                    st.rerun()
            with c3:
                if st.button("← Back", key="retry_step3", use_container_width=True):
                    st.session_state.step = 2
                    st.rerun()
    
    render_section_card("Executive Summary Generation", step3_content)

# ==============================
# STEP 4: Flowchart Generation
# ==============================
elif st.session_state.step == 4:
    def step4_content():
        _ensure_flow_state()
        fs = st.session_state.flow_state

        st.markdown("#### Flowchart Generation")

        # 1) Generate Mermaid from solution PDF
        gen = st.button("Generate flowchart from solution PDF", key="btn_gen_flow_from_pdf")
        if gen:
            with st.spinner("Reading PDF and drafting flowchart…"):
                pdf_bytes = st.session_state.get("solution_pdf_bytes")
                if not pdf_bytes and os.path.exists("temp_uploaded.pdf"):
                    pdf_bytes = open("temp_uploaded.pdf","rb").read()

                if not pdf_bytes:
                    st.error("Please upload the Solution Design PDF in Step 1 (or here above).")
                else:
                    try:
                        txt = extract_relevant_text(pdf_bytes)
                    except Exception:
                        txt = ""
                    code = ""
                    if txt:
                        try:
                            code = ask_llm_mermaid(load_groq(), txt)
                        except Exception:
                            code = ""
                    if not code or not code.lower().startswith("flowchart"):
                        st.error("The LLM didn’t return valid Mermaid. Please try again.")
                    else:
                        st.session_state["mermaid_code"] = code.strip()
                        st.success("Flowchart generated from PDF.")

        # 2) Two columns: Left = Mermaid code | Right = Preview
        PREVIEW_HEIGHT = 520
        col_code, col_diagram = st.columns([1, 1], gap="large")

        # --- LEFT: CODE EDITOR (no 'Update Preview' button; preview updates on edit) ---
        with col_code:
            st.subheader("Mermaid Code")
            mermaid_text = st.text_area(
                "Editable Mermaid",
                value=st.session_state.get("mermaid_code", ""),
                key="mermaid_code_editor",
                height=PREVIEW_HEIGHT,
                help="Copy this code and paste in mermaid.live for editing."
            )
            # keep session in sync
            if mermaid_text != st.session_state.get("mermaid_code", ""):
                st.session_state["mermaid_code"] = mermaid_text
            # Actions aligned to original theme
            st.text(" ")
            ac1, ac2, ac3 = st.columns([1, 2, 1])
            with ac1:
                if st.button("Edit", key="edit_in_drawio_btn", use_container_width=True):
                    try:
                        # Clean start
                        if fs.driver:
                            try: fs.driver.quit()
                            except Exception: pass
                        fs.driver = _DrawIO.open_and_insert(
                            st.session_state.get("mermaid_code",""),
                            headless=False,                 # ALWAYS non-headless
                            download_dir=fs.download_dir
                        )
                        fs.mode = "drawio"
                        fs.last_refresh_ts = time.time()
                        st.success("diagrams.net opened. Make edits there, then click Refresh.")
                    except Exception as e:
                        st.error(f"Failed to open Draw.io: {e}")

        # --- RIGHT: LIVE PREVIEW + ACTIONS (Edit + Refresh) ---
        with col_diagram:
            st.subheader("Preview")
             # Show last refresh time if available
            if fs.last_refresh_ts:
                elapsed = int(time.time() - fs.last_refresh_ts)
                st.caption(f"Last refreshed: {elapsed}s ago")
            if fs.viewer_url and fs.mode == "drawio":
                components.iframe(fs.viewer_url, height=PREVIEW_HEIGHT, scrolling=True)

            else:
                code_for_preview = st.session_state.get("mermaid_code", "").strip()
                if code_for_preview:
                    safe = sanitize_mermaid_for_render(code_for_preview)
                    render_mermaid_chart(safe, height=PREVIEW_HEIGHT + 20)
                    st.caption("Preview rendered in browser with Mermaid.js.")
                else:
                    st.info("No preview yet. Generate from PDF or paste Mermaid code.")

            # Actions aligned to original theme
            # --- replace your refresh button block with this ---

            ac1, ac2, ac3 = st.columns([1, 2, 1])

            with ac3:
                # Make the refresh look like an icon (no blue background), scoped to this column
                st.markdown("""
                    <style>
                        /* Scope to the LAST column of this 3-col row */
                        div[data-testid="column"]:nth-of-type(3) .stButton > button {
                            background: transparent !important;
                            border: none !important;
                            box-shadow: none !important;
                            color: #060c71 !important;            /* brand blue */
                            padding: .25rem .4rem !important;
                            font-size: 32px !important;           /* big icon */
                            line-height: 1 !important;
                            border-radius: 100% !important;
                            transition: transform .15s ease;
                        }
                       
                    </style>
                """, unsafe_allow_html=True)
                # Debounce: prevent rapid clicking
                can_refresh = True
                if fs.last_refresh_ts:
                    elapsed = time.time() - fs.last_refresh_ts
                    if elapsed < 3:  # 3-second cooldown
                        can_refresh = False
                        st.caption(f"⏳ Wait {3-int(elapsed)}s")

                if st.button("↻", key="refresh_now_btn", help="Refresh preview from diagrams.net"):
                    with st.spinner("Exporting...", show_time=True):
                        success = _auto_refresh_drawio_preview()
                        if success:
                            st.rerun()
                            try:
                                st.toast("✅ Preview refreshed", icon="✅")
                            except Exception:
                                pass


        # 3) Next → capture the LATEST flowchart into a 1-page DOCX with header/footer
        st.markdown("---")
        merge_disabled = not all([
            st.session_state.get("selected_cover_buffer"),
            st.session_state.get("selected_summary_buffer"),
            bool(st.session_state.get("mermaid_code","").strip()),
        ])
        
        if merge_disabled:
            missing = []
            if not st.session_state.get("selected_cover_buffer"):   missing.append("Cover Letter")
            if not st.session_state.get("selected_summary_buffer"): missing.append("Executive Summary")
            if not st.session_state.get("mermaid_code","").strip(): missing.append("Flowchart")
            st.warning(f"Missing components for final merge: {', '.join(missing)}")

        c1, c2, c3 = st.columns([1, 2, 1])
        with c1:
            if st.button("Next →", disabled=merge_disabled, key="merge_final_btn", use_container_width=True):
                try:
                    png_bytes = None

                    # 1) If editing in diagrams.net, capture latest view
                    if fs.mode == "drawio" and fs.driver is not None:
                        try:
                            # Export fresh XML → regenerate viewer URL
                            xml_file = _DrawIO.export_xml(fs.driver, fs.download_dir, timeout_sec=45)
                            xml_text = pathlib.Path(xml_file).read_text(encoding="utf-8", errors="ignore")
                            fs.viewer_url = _DrawIO.xml_to_viewer_url(xml_text)
                            
                            # Screenshot the viewer URL
                            if fs.viewer_url:
                                drv = _SeleniumHelper.create_driver(headless=True, download_dir=None)
                                try:
                                    drv.set_window_size(1920, 1200)
                                    drv.get(fs.viewer_url)
                                    time.sleep(2.5)  # Increased wait time
                                    png_bytes = drv.get_screenshot_as_png()
                                finally:
                                    try:
                                        drv.quit()
                                    except Exception:
                                        pass
                        except Exception as e:
                            st.warning(f"Draw.io export failed: {e}, falling back to Mermaid render")

                    # 2) Fallback to rendering current Mermaid code
                    if not png_bytes or len(png_bytes) == 0:
                        mermaid_code = st.session_state.get("mermaid_code", "")
                        if mermaid_code.strip():
                            png_bytes = mermaid_to_png_via_kroki(mermaid_code)
                            if not png_bytes:
                                # Try Chrome fallback
                                png_bytes = mermaid_to_png_via_chrome(mermaid_code)

                    # 3) Verify we have valid PNG bytes
                    if not png_bytes or len(png_bytes) == 0:
                        st.error("⚠️ Could not generate flowchart image. Please verify your flowchart is valid and try again.")
                        st.info("Debug info: Check if mermaid.live can render your code")
                        st.stop()

                    # Store the PNG bytes for reference
                    st.session_state["flowchart_final_png_bytes"] = png_bytes
                    st.session_state["mermaid_png_bytes"] = png_bytes

                    if not png_bytes or len(png_bytes) < 100:
                        st.error("❌ Failed to capture flowchart image. Please try refreshing again.")
                        st.stop()
                    
                    # Verify PNG is valid before proceeding
                    try:
                        test_img = PILImage.open(io.BytesIO(png_bytes))
                        test_img.load()
                        st.info(f"✅ Captured PNG: {test_img.width}x{test_img.height} pixels")
                    except Exception as e:
                        st.error(f"❌ PNG validation failed: {e}")
                        st.stop()

                    # Build the DOCX page with the image
                    flow_docx = make_docx_fit_one_page_from_png_or_code(
                        png_bytes,
                        st.session_state.get("mermaid_code", ""),
                        title="Concept Description",
                        client_logo=st.session_state.get("client_logo"),
                        falcon_logo=st.session_state.get("falcon_logo"),
                        client_name=st.session_state.get("client_name", ""),
                        project_name=st.session_state.get("project_title", ""),
                    )

                    # Verify the DOCX was created
                    if not flow_docx or flow_docx.getbuffer().nbytes < 5000:
                        st.error("Failed to create flowchart document page")
                        st.stop()

                    # Store the buffer for final merge
                    st.session_state.flowchart_docx_buffer = flow_docx
                    st.session_state["flowchart_is_final"] = True

                    st.success("✅ Flowchart page prepared successfully!")
                    time.sleep(0.5)  # Brief pause to show success message
                    st.session_state.step = 5
                    st.rerun()

                except Exception as e:
                    st.error(f"❌ Error preparing flowchart: {str(e)}")
                    import traceback
                    st.code(traceback.format_exc())
                    st.stop()

        with c3:
            if st.button("← Back", key="back_step4", use_container_width=True):
                st.session_state.step = 3
                st.rerun()

    render_section_card("Flowchart Generation", step4_content)


# ==============================
# STEP 5: Final Proposal (Merge & Download)
# ==============================
elif st.session_state.step == 5:
    def step5_content():
        st.caption("Merge Everything into a single DOCX.")

        # --- Build cover page using YOUR create_cover_page() ---
        def build_cover_from_template() -> io.BytesIO:
            # You can set a template path earlier via st.session_state.cover_template_path
            template_path = st.session_state.get("cover_template_path", "Input\\Cover_Temp.docx")
            client_logo = st.session_state.get("client_logo")
            client_name = st.session_state.get("client_name", "Client")
            project_title = st.session_state.get("project_title", "Project")

            if os.path.exists(template_path):
                try:
                    return create_cover_page(
                        template_path=template_path,
                        client_logo=client_logo,
                        client_name=client_name,
                        project_title=project_title,
                    )
                except Exception as e:
                    st.warning(f"Cover template found but failed to render, falling back to simple cover. Error: {e}")

            # ----- Fallback: simple cover -----
            doc = Document()
            # logo
            if client_logo:
                try:
                    p_logo = doc.add_paragraph()
                    p_logo.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    run_logo = p_logo.add_run()
                    run_logo.add_picture(io.BytesIO(client_logo), width=Inches(2.0))
                except Exception:
                    pass
            # spacing
            for _ in range(3):
                doc.add_paragraph("")
            # title/date
            p_title = doc.add_paragraph()
            r_title = p_title.add_run(f"FALCON’s Proposal to {client_name} for the {project_title}")
            r_title.font.name = "Calibri"; r_title.font.size = Pt(24); r_title.bold = True
            p_date = doc.add_paragraph()
            r_date = p_date.add_run(datetime.today().strftime("%B %d, %Y"))
            r_date.font.name = "Calibri"; r_date.font.size = Pt(14)

            buf = io.BytesIO(); doc.save(buf); buf.seek(0)
            return buf

        def build_company_profile_buffer() -> io.BytesIO:
            doc = Document()
            # Branded headers/footers across company profile pages
            add_page_headers(
                doc,
                st.session_state.get("client_logo"),
                st.session_state.get("falcon_logo"),
                st.session_state.get("client_name",""),
                st.session_state.get("project_title",""),
            )
            build_company_profile(doc)  # your existing function fills all content/images
            out = io.BytesIO(); doc.save(out); out.seek(0)
            return out

        # Ensure flowchart page docx exists and is up-to-date with header/footer
        if not st.session_state.get("flowchart_docx_buffer"):
            # === ALWAYS rebuild the flowchart page from the latest diagram before merge ===
            try:
                # Prefer a fresh export if diagrams.net editor is active
                latest_png = _export_latest_flowchart_png()
            except Exception:
                latest_png = None

            if not latest_png:
                # fallbacks: whatever we captured at Step 4, then any previous render
                latest_png = (st.session_state.get("flowchart_final_png_bytes")
                              or st.session_state.get("mermaid_png_bytes"))

            if latest_png and len(latest_png) > 0:
                try:
                    st.session_state.flowchart_docx_buffer = make_docx_fit_one_page_from_png_or_code(
                        latest_png,
                        st.session_state.get("mermaid_code", ""),
                        title="Concept Description",
                        client_logo=st.session_state.get("client_logo"),
                        falcon_logo=st.session_state.get("falcon_logo"),
                        client_name=st.session_state.get("client_name", ""),
                        project_name=st.session_state.get("project_title", ""),
                    )
                except Exception as e:
                    st.error(f"Could not rebuild flowchart page: {e}")
                    st.stop()
            else:
                st.error("⚠️ No flowchart image is available to merge. Please go back to Step 4 and refresh/export the diagram.")
                st.stop()

        # --- Gather parts in required order ---
        parts: List[io.BytesIO] = []
        labels: List[str] = []

        # 1) Cover Page (always built)
        cover_page = build_cover_from_template()
        parts.append(cover_page); labels.append("Cover Page")

        # 2) Cover Letter (selected)
        cov = st.session_state.get("selected_cover_buffer")
        if cov:
            parts.append(cov); labels.append("Cover Letter")

        # 3) Response to RFQ (diagram page you built in Step 1)
        rfq = st.session_state.get("manual_rfq_buffer")
        if rfq:
            parts.append(rfq); labels.append("Response to RFQ")

        # 4) Executive Summary (selected)
        es = st.session_state.get("selected_summary_buffer")
        if es:
            es = _indent_exec_summary_bullets(es, one_indent_in=0.25, base_left_in=0.5, hang_in=0.25)
            parts.append(es); labels.append("Executive Summary")

        # 5) Company Profile (always built)
        company_profile = build_company_profile_buffer()
        parts.append(company_profile); labels.append("Company Profile")

        # 6) Concept Description (latest flowchart one-pager with headers/footers)
        flow = st.session_state.get("flowchart_docx_buffer")
        if flow:
            parts.append(flow); labels.append("Concept Description")

        # ===============================
        # 📋 UI: Section availability checklist (tick/cross)
        # ===============================
        # Derive dynamic availability for the requested list
        def yesno(b): return "✅" if b else "❌"
        section_status = [
            ("Cover page", True),                                  # built above
            ("Cover letter", bool(cov)),
            ("Response to RFQ", bool(rfq)),
            ("Table of content", False),
            ("Glossary", False),
            ("Executive summary", bool(es)),
            ("Company profile", True),                             # built above
            ("News headlines", True),
            ("Reference project", False),
            ("Handled shipment spectrum", False),
            ("Proposed system description", False),
            ("Concept description (Flowchart)", bool(flow)),
            ("Capacity calculation", False),
            ("System Description", False),
            ("Component equipment description", False),
            ("Technical specifications", False),
            ("BOM details", False),
            ("electrical details", False),
            ("Key componenets make", False),
            ("Program organizaion", False),
            ("Commercial terms and condition", False),
            ("warranty and exclusion", False),
        ]

        st.markdown("##### Sections status in this build")
        left, right = st.columns(2)
        half = (len(section_status) + 1) // 2
        for i, (name, present) in enumerate(section_status):
            tgt = left if i < half else right
            with tgt:
                st.markdown(f"{yesno(present)} **{name}**")

        # Validate required items
        missing = []
        if not cov:  missing.append("Cover Letter")
        if not rfq:  missing.append("Response to RFQ")
        if not es:   missing.append("Executive Summary")
        if not flow: missing.append("Concept Description")

        st.write("**Will include (actual merge order):** " + (", ".join(labels) if labels else "—"))
        if missing:
            st.warning("Missing required sections: " + ", ".join(missing))
        can_merge = all([cov, rfq, es, flow])

        colA, colB = st.columns([1, 1])
        with colA:
            gen_btn = st.button(
                "Generate Final Proposal",
                disabled=not can_merge,
                use_container_width=True,
                key="btn_generate_final_proposal_all"
            )

        if gen_btn:
            try:
                merged_buffer = merge_docx_files_with_page_breaks(parts)
                st.session_state["final_docx_bytes"] = merged_buffer.getvalue()

                st.success("Final proposal created successfully!")

                fname = f'{st.session_state.get("client_name","Client")}_{st.session_state.get("project_title","Proposal")}_Final_Proposal.docx'
                st.download_button(
                    label="Download Final Proposal",
                    data=st.session_state["final_docx_bytes"],
                    file_name=fname,
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True,
                    key="btn_download_final_proposal_all"
                )
            except Exception as e:
                st.error(f"Final merge failed: {e}")

        st.markdown("---")
        nav1, nav2, nav3 = st.columns([1, 2, 1])
        with nav3:
            if st.button("← Back", key="back_to_flowchart", use_container_width=True):
                st.session_state.step = 4
                st.rerun()
        with nav1:
            if st.button("Restart ↺", key="restart_wizard", type="primary", use_container_width=True):
                st.session_state.step = 1
                st.rerun()

    render_section_card("Final Proposal", step5_content)
